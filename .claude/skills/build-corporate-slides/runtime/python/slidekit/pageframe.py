"""ページ外枠（表紙・章扉・通常ページのヘッダー）を組み立てる。

いずれもcontent_region()の内側に置く「部品」ではなく、スライド1枚の
外枠そのものを描く。呼び出し元はDeckBuilder（builder.py）だけで、
renderer層はDeckBuilder経由でのみ使う。ページ内の部品はcomponents.py。
層構成の全体像はARCHITECTURE.mdを参照。
"""
from datetime import date

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .atoms import Marker, _flat, add_hairline
from .theme import LAYOUT, PALETTE
from .typography import Tag, _type_for, add_textbox


def add_slide_title(slide, title, *, kicker=None, page=None):
    typography = _type_for(slide)
    if kicker:
        add_textbox(slide, LAYOUT.margin_x, Inches(0.34), Inches(5.5), Inches(0.2),
                    kicker.upper(), size=typography.small, color=PALETTE.blue, bold=True)
    add_textbox(slide, LAYOUT.margin_x, Inches(0.62), Inches(11.7), Inches(0.52),
                title, size=typography.title, color=PALETTE.ink, bold=True,
                font=typography.headline_font)
    add_hairline(slide, LAYOUT.margin_x, Inches(1.25),
                 LAYOUT.slide_width - 2 * LAYOUT.margin_x,
                 color=PALETTE.line_brand, width=1.4)
    if page is not None:
        add_textbox(slide, Inches(12.1), Inches(0.37), Inches(0.55), Inches(0.2),
                    f"{page:02}", size=typography.small, color=PALETTE.grey_500)


def add_cover(slide, title, *, subtitle=None, department="", created=None,
              eyebrow=None, logo_path=None, classification="",
              brand_side="right", brand_shape="diagonal", brand_width=None):
    """表紙を描く。brand_widthでブランド面の幅を調整できる（既定は3.4in/2.84in）。

    左配置時はcontent_xをbrand_widthに応じて自動的に押し出す。右配置時に
    大きく広げる場合は、本文が8.5in幅に収まっているか確認する。
    """
    typography = _type_for(slide)
    created = created or date.today().isoformat()
    if brand_side not in {"left", "right", "none"}:
        raise ValueError("brand_sideは 'left'、'right'、'none' を指定してください")
    if brand_shape not in {"diagonal", "curve", "straight"}:
        raise ValueError("brand_shapeは 'diagonal'、'curve'、'straight' を指定してください")

    default_widths = {"curve": Inches(3.4), "straight": Inches(2.84),
                      "diagonal": Inches(3.4)}
    width = brand_width or default_widths[brand_shape]

    if brand_side != "none":
        brand_shapes = []
        if brand_shape == "curve":
            inner_w = int(width * (1.25 / 3.4))
            if brand_side == "left":
                brand_shapes.append(slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, inner_w, LAYOUT.slide_height))
                brand_shapes.append(slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, 0, 0, width, LAYOUT.slide_height))
            else:
                brand_shapes.append(slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, LAYOUT.slide_width - inner_w, 0,
                    inner_w, LAYOUT.slide_height))
                brand_shapes.append(slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, LAYOUT.slide_width - width, 0,
                    width, LAYOUT.slide_height))
        elif brand_shape == "straight":
            x = 0 if brand_side == "left" else LAYOUT.slide_width - width
            brand_shapes.append(slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, 0, width, LAYOUT.slide_height))
        else:
            # 斜めの境界を一つだけにし、片側全体を単色の面にする
            # （PARALLELOGRAMは両辺が斜めになり境界が二重に見えるため使わない）。
            # 下側（ロゴを置く側）をやや広く、上側（タイトル側）をやや狭く
            # 取り、単調な平行四辺形にならないようにする。
            narrow_w = int(width * 0.76)
            wide_w = int(width * 1.18)
            if brand_side == "left":
                points = [(0, 0), (narrow_w, 0),
                         (wide_w, LAYOUT.slide_height),
                         (0, LAYOUT.slide_height)]
            else:
                top_x = LAYOUT.slide_width - narrow_w
                bottom_x = LAYOUT.slide_width - wide_w
                points = [(top_x, 0), (LAYOUT.slide_width, 0),
                         (LAYOUT.slide_width, LAYOUT.slide_height),
                         (bottom_x, LAYOUT.slide_height)]
            builder = slide.shapes.build_freeform(
                start_x=points[0][0], start_y=points[0][1])
            builder.add_line_segments(points[1:], close=True)
            brand_shapes.append(builder.convert_to_shape())
        for brand in brand_shapes:
            _flat(brand)
            brand.line.fill.background()
            if brand_shape == "diagonal":
                # 単色の面だけだと平板に見えるため、ごく控えめなグラデー
                # ションで奥行きを出す（表紙のブランド面に限定した使用）。
                brand.fill.gradient()
                brand.fill.gradient_angle = 115.0
                stops = brand.fill.gradient_stops
                stops[0].position = 0.0
                stops[0].color.rgb = RGBColor(0xE7, 0xF1, 0xFD)
                stops[1].position = 1.0
                stops[1].color.rgb = RGBColor(0xCC, 0xDB, 0xF8)
            else:
                brand.fill.solid()
                brand.fill.fore_color.rgb = PALETTE.surface_brand_soft

    if brand_side == "left":
        content_x = max(Inches(3.25), width + Inches(0.35))
        content_w = LAYOUT.slide_width - LAYOUT.margin_x - content_x
    else:
        content_x = LAYOUT.margin_x
        content_w = Inches(8.5)
    if eyebrow:
        add_textbox(slide, content_x, Inches(2.08), content_w, Inches(0.28),
                    eyebrow.upper(), size=typography.small, color=PALETTE.blue, bold=True)
    # 表紙はポスター的な大見出しのため、本文12pt基準で縮むtypography.titleは
    # 使わず、add_section_dividerと同じ固定サイズにする。
    add_textbox(slide, content_x, Inches(2.48), content_w, Inches(1.3),
                title, size=Pt(34), color=PALETTE.text_primary, bold=True,
                font=typography.headline_font, line_spacing=0.98)
    if subtitle:
        add_textbox(slide, content_x, Inches(4.03), content_w, Inches(0.62),
                    subtitle, size=typography.section, color=PALETTE.text_secondary)

    meta_x, meta_w = Inches(9.2), Inches(3.45)
    badge_x, badge_y, badge_w = Inches(11.42), Inches(0.32), Inches(1.22)
    if classification:
        # 開示区分バッジは社内文書の型として角ありで固定運用する
        # （実務要件。他の一般的なTagはpill=Trueの丸薬型のまま）。
        Tag(slide, badge_x, badge_y, badge_w, Inches(0.3), classification,
           pill=False)
    if department:
        add_textbox(slide, meta_x, Inches(0.72), meta_w, Inches(0.24), department,
                    size=typography.small, color=PALETTE.text_secondary,
                    bold=True, align=PP_ALIGN.RIGHT)
    add_textbox(slide, meta_x, Inches(1.0), meta_w, Inches(0.24), created,
                size=typography.small, color=PALETTE.text_secondary,
                align=PP_ALIGN.RIGHT)
    if logo_path:
        # 実際の右下の角へ、部外秘ボックスの右端・上端の余白を基準に置く
        # （同じ余白量で角に揃えると、意図的な配置に見える）。
        logo_w = Inches(0.55)
        corner_margin = badge_y  # 部外秘ボックスの上端の余白と揃える
        logo_x = Inches(0.55) if brand_side == "left" else badge_x + badge_w - logo_w
        picture = slide.shapes.add_picture(str(logo_path), logo_x, 0, width=logo_w)
        picture.top = LAYOUT.slide_height - corner_margin - picture.height


def add_section_divider(slide, title, *, kicker=None, subtitle=None, page=None):
    """章扉。通常ページの小見出しヘッダーは使わず、単独で成立させる。"""
    typography = _type_for(slide)
    center_y = Inches(3.2)
    Marker(slide, LAYOUT.margin_x, center_y, Inches(0.09), Inches(0.6))
    text_x = LAYOUT.margin_x + Inches(0.3)
    title_y = center_y
    if kicker:
        add_textbox(slide, text_x, center_y - Inches(0.02), Inches(10),
                    Inches(0.3), kicker.upper(), size=typography.small,
                    color=PALETTE.blue, bold=True)
        title_y = center_y + Inches(0.36)
    add_textbox(slide, text_x, title_y, Inches(11), Inches(1.2), title,
                size=Pt(34), color=PALETTE.text_primary, bold=True,
                font=typography.headline_font, line_spacing=1.02)
    if subtitle:
        add_textbox(slide, text_x, title_y + Inches(1.05), Inches(9.5),
                    Inches(0.6), subtitle, size=typography.section,
                    color=PALETTE.text_secondary)
    if page is not None:
        add_textbox(slide, Inches(12.1), Inches(0.37), Inches(0.55),
                    Inches(0.2), f"{page:02}", size=typography.small,
                    color=PALETTE.grey_500)
