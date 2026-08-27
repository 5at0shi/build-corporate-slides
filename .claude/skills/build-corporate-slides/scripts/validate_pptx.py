#!/usr/bin/env python3
import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "runtime" / "python"))
from slidekit.textmetrics import char_width_factor  # noqa: E402


_THIN_W = Inches(0.12)
_THIN_H = Inches(0.16)


def _bbox(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _overlap_area(a, b):
    left, top = max(a[0], b[0]), max(a[1], b[1])
    right, bottom = min(a[2], b[2]), min(a[3], b[3])
    if right <= left or bottom <= top:
        return 0
    return (right - left) * (bottom - top)


def _text_glyph_bbox(shape):
    """テキストが実際に占めるであろう概算の矩形。

    textboxは中身より高さに余裕を持たせて作ることが多い（例:
    add_text_list(divider=True)の各行は行の高さに関わらず常に0.76in）。図形の公称の
    bboxをそのまま重なり判定に使うと、余白部分に触れただけの装飾を
    誤検知してしまうため、vertical_anchorと推定文字高さから実際の
    文字が占める範囲だけを切り出す。
    """
    box = _bbox(shape)
    tf = shape.text_frame
    estimated_pt = _estimate_text_height_pt(tf, Emu(shape.width).pt)
    if estimated_pt is None:
        return box
    estimated_emu = int(estimated_pt * 12700)
    top, height = shape.top, shape.height
    if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
        glyph_top = top + max(0, (height - estimated_emu) // 2)
    elif tf.vertical_anchor == MSO_ANCHOR.BOTTOM:
        glyph_top = top + max(0, height - estimated_emu)
    else:
        glyph_top = top
    glyph_bottom = min(top + height, glyph_top + estimated_emu)
    return (box[0], glyph_top, box[2], glyph_bottom)


def _check_decoration_overlap(slide, slide_no, warnings):
    """罫線・アクセントバー等の薄い装飾図形が、文字の上に被っていないか検知する。

    add_card/add_background_zone等の大きな背景面はテキストを内側に持つのが
    通常の使い方（薄くない図形は対象外）なので誤検知しない。角丸統一・
    org_layersの見出しマーカー・numbered_listの罫線で、実際にこの種の
    重なりが起きたことがあるため、目視に頼らず機械的に検知する。
    """
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            continue
        if shape.width <= 0 or shape.height <= 0:
            continue
        if shape.width > _THIN_W and shape.height > _THIN_H:
            continue
        deco_area = shape.width * shape.height
        deco_box = _bbox(shape)
        for text_shape in text_shapes:
            overlap = _overlap_area(deco_box, _text_glyph_bbox(text_shape))
            if overlap / deco_area > 0.12:
                snippet = text_shape.text.strip().replace("\n", " ")[:24]
                warnings.append(
                    f"slide {slide_no}: 薄い装飾図形 '{shape.name}' が"
                    f"テキスト '{snippet}' に重なっている可能性")
                break


def _paragraph_line_sizes(paragraph, available_w: float) -> list[float] | None:
    """runをまたいだ折り返しを再現し、各行に使われた最大フォントサイズを返す。

    add_text_list（item_list/numbered_list/icon_list）などは、タイトルと本文を改行文字("\\n")
    で1つのparagraph内の別runへ連結する（見た目上の行ごとにtextboxを
    分割しないため）。単純に全runのテキストを結合して1つのフォントサイズで
    折り返し推定すると、明示的な改行を無視して過小評価する。runごとに
    改行で区切り、行に使われた文字の最大フォントサイズで高さを見積もる。
    """
    sizes = [run.font.size.pt for run in paragraph.runs if run.font.size]
    if not sizes:
        return None
    fallback_size = max(sizes)
    line_sizes = []
    current_w = 0.0
    current_max_size = 0.0
    started = False
    for run in paragraph.runs:
        size = run.font.size.pt if run.font.size else fallback_size
        segments = run.text.split("\n")
        for seg_index, segment in enumerate(segments):
            if seg_index > 0:
                line_sizes.append(current_max_size or size)
                current_w = 0.0
                current_max_size = 0.0
            for ch in segment:
                started = True
                ch_w = char_width_factor(ch) * size
                if current_w + ch_w > available_w and current_w > 0:
                    line_sizes.append(current_max_size or size)
                    current_w = 0.0
                    current_max_size = 0.0
                current_w += ch_w
                current_max_size = max(current_max_size, size)
    if current_w > 0 or not line_sizes:
        line_sizes.append(current_max_size or fallback_size)
    return line_sizes if started else None


def _estimate_text_height_pt(text_frame, width_pt: float) -> float | None:
    """textboxの折り返し後の推定高さ(pt)。フォントサイズ未指定の場合はNone。

    生成時のレイアウト判断(slidekit.textmetrics)と同じ文字幅ヒューリス
    ティックを使い、runをまたいだ改行も再現して折り返し行数を見積もる。
    """
    margin_left = Emu(text_frame.margin_left or 0).pt
    margin_right = Emu(text_frame.margin_right or 0).pt
    margin_top = Emu(text_frame.margin_top or 0).pt
    margin_bottom = Emu(text_frame.margin_bottom or 0).pt
    available_w = max(1.0, width_pt - margin_left - margin_right)
    total = margin_top + margin_bottom
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs or not any(run.text for run in paragraph.runs):
            continue
        line_sizes = _paragraph_line_sizes(paragraph, available_w)
        if line_sizes is None:
            return None
        spacing = paragraph.line_spacing
        multiplier = spacing if isinstance(spacing, (int, float)) else 1.0
        space_before = Emu(paragraph.space_before or 0).pt
        space_after = Emu(paragraph.space_after or 0).pt
        total += sum(size * 1.2 * multiplier for size in line_sizes)
        total += space_before + space_after
    return total


def _check_table(shape, slide_no, issues, warnings):
    """表はGraphicFrameでhas_text_frame=Falseのため、他のテキストチェックの
    対象から漏れる。セルの小さすぎる文字と、行の高さに対する文字のはみ出し
    をここで別途確認する。
    """
    table = shape.table
    # 行ごとに「最も高さを要するセル」を求め、表全体の実際の高さを見積もる。
    # PowerPointは行の高さを内容に合わせて自動的に広げる（縮めない）ため、
    # python-pptxが持つ宣言上の高さが枠内でも、描画時には下へ伸びて後続の
    # 行がページ外へ押し出され、そのまま見えなくなる。宣言値だけを見る
    # スライド外チェックでは検出できないので、ここで推定して確かめる。
    declared_pt = 0.0
    estimated_total_pt = 0.0
    for row_index in range(len(table.rows)):
        row_height_pt = Emu(table.rows[row_index].height).pt
        declared_pt += row_height_pt
        tallest_pt = row_height_pt
        for col_index in range(len(table.columns)):
            cell = table.cell(row_index, col_index)
            text = cell.text_frame.text
            if not text.strip():
                continue
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size.pt < 7:
                        issues.append(
                            f"slide {slide_no}: 表内に7pt未満の文字 '{run.text[:24]}'")
            col_width_pt = Emu(table.columns[col_index].width).pt
            estimated = _estimate_text_height_pt(cell.text_frame, col_width_pt)
            if estimated is None:
                continue
            tallest_pt = max(tallest_pt, estimated)
            if estimated > row_height_pt * 1.15:
                snippet = text.strip().replace("\n", " ")[:24]
                warnings.append(
                    f"slide {slide_no}: 表のセルが行の高さからはみ出す可能性 "
                    f"'{snippet}' (推定{estimated:.0f}pt / 行{row_height_pt:.0f}pt)")
        estimated_total_pt += tallest_pt

    # 表全体が確保した高さを超えると、はみ出した分だけ後ろの行が下へ押し出され、
    # ページ下端の結論やスライド外へ隠れて読めなくなる（データが消える）。
    # 見た目が窮屈になるだけの個別セルの警告とは重大度が違うためissueにする。
    if declared_pt > 0 and estimated_total_pt > declared_pt * 1.05:
        overflow_pt = estimated_total_pt - declared_pt
        lost_rows = max(1, int(overflow_pt / (declared_pt / len(table.rows))))
        issues.append(
            f"slide {slide_no}: 表が確保した高さに収まりません "
            f"(推定{estimated_total_pt:.0f}pt / 確保{declared_pt:.0f}pt)。"
            f"末尾およそ{lost_rows}行がページ下端へ押し出されて見えなくなります。"
            "行数を減らすか、セルの文言を短くするか、ページを分割してください")


def _check_text_overflow(slide, slide_no, warnings):
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text.strip():
            continue
        if shape.height <= 0 or shape.width <= 0:
            continue
        estimated = _estimate_text_height_pt(shape.text_frame, Emu(shape.width).pt)
        if estimated is None:
            continue
        actual = Emu(shape.height).pt
        if estimated > actual * 1.15:
            snippet = shape.text.strip().replace("\n", " ")[:24]
            warnings.append(
                f"slide {slide_no}: テキストが枠からはみ出す可能性 "
                f"'{snippet}' (推定{estimated:.0f}pt / 枠{actual:.0f}pt)")


def validate(path: Path) -> tuple[list[str], list[str]]:
    issues = []
    warnings = []
    if not path.is_file() or path.stat().st_size == 0:
        return [f"PPTXが存在しないか空です: {path}"], warnings
    try:
        prs = Presentation(path)
    except Exception as exc:
        return [f"PPTXを開けません: {exc}"], warnings
    if not prs.slides:
        issues.append("スライドがありません")
    for slide_no, slide in enumerate(prs.slides, 1):
        visible = 0
        text_shapes = []
        for shape in slide.shapes:
            is_line = shape.shape_type == MSO_SHAPE_TYPE.LINE
            invalid_size = ((shape.width <= 0 and shape.height <= 0) if is_line
                            else (shape.width <= 0 or shape.height <= 0))
            if invalid_size:
                issues.append(f"slide {slide_no}: サイズ0のオブジェクト '{shape.name}'")
            if (shape.left + shape.width < 0 or shape.top + shape.height < 0 or
                    shape.left > prs.slide_width or shape.top > prs.slide_height):
                issues.append(f"slide {slide_no}: スライド外のオブジェクト '{shape.name}'")
            if shape.has_text_frame and shape.text.strip():
                visible += 1
                text_shapes.append(shape)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size.pt < 7:
                            issues.append(
                                f"slide {slide_no}: 7pt未満の文字 '{run.text[:24]}'")
            elif not shape.has_text_frame:
                visible += 1
                if shape.has_table:
                    _check_table(shape, slide_no, issues, warnings)
        if visible == 0:
            issues.append(f"slide {slide_no}: 空ページの可能性")
        _check_text_overflow(slide, slide_no, warnings)
        _check_decoration_overlap(slide, slide_no, warnings)
        short_shapes = [s for s in text_shapes if len(s.text.strip()) <= 24]
        if len(text_shapes) >= 16 and len(short_shapes) / len(text_shapes) >= 0.65:
            warnings.append(
                f"slide {slide_no}: 短いtextboxが多く、文章の過剰分割の可能性 "
                f"({len(text_shapes)} text shapes)")
    return issues, warnings


def main() -> int:
    parser = argparse.ArgumentParser(description="PPTXの基本構造を検証します")
    parser.add_argument("pptx", type=Path)
    args = parser.parse_args()
    issues, warnings = validate(args.pptx)
    if warnings:
        print("編集性の警告:")
        for warning in warnings:
            print(f"- {warning}")
    if issues:
        print("検証で問題を検出しました:")
        for issue in issues:
            print(f"- {issue}")
        return 1
    print(f"OK: {args.pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
