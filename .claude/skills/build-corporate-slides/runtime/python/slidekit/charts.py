from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION,
                             XL_LEGEND_POSITION, XL_MARKER_STYLE)
from pptx.util import Pt

from .theme import PALETTE

_CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}
_SCATTER_TYPE = XL_CHART_TYPE.XY_SCATTER

# 系列・要素の配色はsemantic paletteを巡回させ、デッキ全体の配色と揃える。
_SERIES_COLORS = [PALETTE.blue, PALETTE.accent_secondary, PALETTE.grey_500,
                  PALETTE.navy]


def add_native_chart(slide, x, y, w, h, *, chart_type, categories=None,
                     series, typography, value_format=None):
    """編集可能なネイティブPowerPointグラフを描く。

    chart_typeは column/stacked_column/bar/line/pie/scatter。
    column/stacked_column/bar/line/pieはcategoriesと
    series=[{"name": ..., "values": [...]}] を組み合わせるカテゴリ型。
    scatterはcategories不要で、
    series=[{"name": ..., "points": [{"x": ..., "y": ...}, ...]}] という
    XY座標形式を使う（ポートフォリオ分析等、マスのどこかではなく正確な
    座標が要点の場合。マスの位置づけだけが要点ならmatrixを使う）。
    PNG画像と異なり、貼り付け後もPowerPoint上で数値・系列名を直接編集できる。
    """
    if chart_type == "scatter":
        return _add_scatter_chart(slide, x, y, w, h, series=series,
                                  typography=typography)
    if chart_type not in _CHART_TYPES:
        raise ValueError(f"未対応のchart_typeです: {chart_type}")
    if not categories:
        raise ValueError("categoriesが空です")
    if not series:
        raise ValueError("seriesが空です")
    if chart_type == "pie" and len(series) > 1:
        raise ValueError(
            "pieは1系列のみ対応です（PowerPointの円グラフは最初の系列しか"
            f"表示せず、他は数値だけ残って混乱のもとになります）: {len(series)}系列")
    for s in series:
        values = s.get("values", [])
        if len(values) != len(categories):
            raise ValueError(
                f"系列'{s.get('name', '')}'のvalues数({len(values)})が"
                f"categories数({len(categories)})と一致しません")

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    graphic_frame = slide.shapes.add_chart(
        _CHART_TYPES[chart_type], x, y, w, h, chart_data)
    chart = graphic_frame.chart
    chart.has_title = False
    font = typography.body_font
    small = typography.small

    has_negative = any(v < 0 for s in series for v in s.get("values", []))
    if chart_type == "pie":
        _style_pie(chart, font, small)
    else:
        _style_category_chart(chart, chart_type, font, small, value_format,
                              multi_series=len(series) > 1,
                              has_negative=has_negative,
                              category_count=len(categories))
    return graphic_frame


def _style_pie(chart, font, small):
    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.show_category_name = True
    labels.show_percentage = True
    labels.show_value = False
    labels.number_format = "0%"
    labels.number_format_is_linked = False
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.font.size = small
    labels.font.name = font
    labels.font.color.rgb = PALETTE.text_primary
    chart.has_legend = False

    points = plot.series[0].points
    for index, point in enumerate(points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = (
            _SERIES_COLORS[index % len(_SERIES_COLORS)])
        point.format.line.color.rgb = PALETTE.surface_base
        point.format.line.width = Pt(1.5)


_MAX_LABELED_CATEGORIES = 8


def _style_category_chart(chart, chart_type, font, small, value_format,
                          *, multi_series, has_negative=False, category_count=0):
    is_stacked = chart_type == "stacked_column"
    plot = chart.plots[0]
    plot.gap_width = 60
    if is_stacked:
        # 積み上げは系列を並べず完全に重ねる。クラスタ型の-8（隙間を作る
        # ための負のoverlap）とは意味が異なる。
        plot.overlap = 100
    elif chart_type != "line":
        plot.overlap = -8 if multi_series else 0

    chart.has_legend = multi_series
    if multi_series:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = small
        chart.legend.font.name = font

    # カテゴリ数が多いと点ごとのラベルが重なって読めなくなるため、閾値を
    # 超えたら軸・目盛線での判読に任せてラベルは付けない。積み上げは
    # セグメントが縦に分かれるため、クラスタ型ほど多系列でも詰まらない。
    plot.has_data_labels = (category_count <= _MAX_LABELED_CATEGORIES
                            and (is_stacked or not multi_series))
    if plot.has_data_labels:
        labels = plot.data_labels
        labels.font.size = small
        labels.font.name = font
        labels.font.color.rgb = PALETTE.text_primary
        if value_format:
            labels.number_format = value_format
            labels.number_format_is_linked = False
        if is_stacked:
            # セグメント内に収める。INSIDE_END/OUTSIDE_ENDは積み上げでは
            # 意味が薄い（隣接セグメントに重なる／棒の外に出て浮く）。
            labels.position = XL_LABEL_POSITION.CENTER
        elif chart_type != "line":
            # 負の値を含む場合、OUTSIDE_ENDだと負の棒のラベルが軸ラベルの行と
            # 重なって読めなくなる（Keynote実測で確認）。棒の内側に収める。
            labels.position = (XL_LABEL_POSITION.INSIDE_END if has_negative
                              else XL_LABEL_POSITION.OUTSIDE_END)

    for index, s in enumerate(plot.series):
        color = _SERIES_COLORS[index % len(_SERIES_COLORS)]
        if chart_type == "line":
            s.format.line.color.rgb = color
            s.format.line.width = Pt(2.25)
            s.smooth = False
        else:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = color
            s.format.line.fill.background()

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = small
    category_axis.tick_labels.font.name = font
    category_axis.format.line.color.rgb = PALETTE.line_neutral

    value_axis = chart.value_axis
    # 棒・列グラフは「棒の長さ」が量を表すため、軸の下限を0に固定する。
    # 自動スケールに任せると、値が0から離れているほど下限が持ち上がり、
    # 例えば[42, 35, 38]で下限30が選ばれて1.2倍の差が3倍に見える。
    # 折れ線・散布図は位置で読む図法なので0起点を強制しない（変化の幅を
    # 見せるために0を含めないことが正当な場面がある）。
    # 負の値を含む場合は下限を触らない。0より下を切り落としてしまうため。
    if chart_type in ("column", "stacked_column", "bar") and not has_negative:
        value_axis.minimum_scale = 0
    value_axis.format.line.fill.background()
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = PALETTE.line_neutral
    value_axis.major_gridlines.format.line.width = Pt(0.5)
    value_axis.tick_labels.font.size = small
    value_axis.tick_labels.font.name = font
    if value_format:
        value_axis.tick_labels.number_format = value_format
        value_axis.tick_labels.number_format_is_linked = False


def _add_scatter_chart(slide, x, y, w, h, *, series, typography):
    """XY座標型の散布図（マーカーのみ、線で結ばない）。

    categoryChartとは異なりcategoriesを使わず、系列ごとに
    points=[{"x": ..., "y": ...}, ...]というXY座標の配列を持つ。
    """
    if not series:
        raise ValueError("seriesが空です")
    chart_data = XyChartData()
    for s in series:
        points = s.get("points", [])
        if not points:
            raise ValueError(f"系列'{s.get('name', '')}'のpointsが空です")
        xy_series = chart_data.add_series(s.get("name", ""))
        for point in points:
            if "x" not in point or "y" not in point:
                raise ValueError(
                    f"系列'{s.get('name', '')}'のpointsにx/yが必要です: {point}")
            xy_series.add_data_point(point["x"], point["y"])

    graphic_frame = slide.shapes.add_chart(
        _SCATTER_TYPE, x, y, w, h, chart_data)
    chart = graphic_frame.chart
    chart.has_title = False
    font = typography.body_font
    small = typography.small

    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = small
        chart.legend.font.name = font

    # python-pptxのXY散布図チャートXMLはhas_data_labels=Falseの明示設定に
    # 対応していない（_remove_dLblsが実装されていない）。新規チャートは
    # 既定でラベル無しのため、明示設定自体を省略する。
    plot = chart.plots[0]
    for index, s in enumerate(plot.series):
        color = _SERIES_COLORS[index % len(_SERIES_COLORS)]
        s.format.line.fill.background()
        s.marker.style = XL_MARKER_STYLE.CIRCLE
        s.marker.size = 8
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = color
        s.marker.format.line.fill.background()

    # scatterは2軸ともvalue axis（category_axisはX軸を指す）。
    for axis in (chart.category_axis, chart.value_axis):
        axis.format.line.color.rgb = PALETTE.line_neutral
        axis.has_major_gridlines = True
        axis.major_gridlines.format.line.color.rgb = PALETTE.line_neutral
        axis.major_gridlines.format.line.width = Pt(0.5)
        axis.tick_labels.font.size = small
        axis.tick_labels.font.name = font
    return graphic_frame
