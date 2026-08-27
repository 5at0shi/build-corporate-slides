"""Atom層のうち、図形を扱う部分（文字はtypography.py、公開している
部品名はcomponents.py）。層構成の全体像はARCHITECTURE.mdを参照。

Boxは「矩形領域＋見た目（塗り・枠線・角丸・影）」という1つの構造を
表す。Card / Background Zone / Focus Panelは、この構造の見た目が違う
だけの同じもの（skinの違い）であり、幾何計算（角丸半径の絶対値換算、
影の作り方）は1箇所に集約する。components.pyのadd_card等は、この
Boxへ薄く委譲するラッパーとして再定義する。

Box/Marker/add_hairlineはいずれも「図形を作る→既定のテーマスタイルを
消す→塗るか透明か→枠線を引くかなしか」という同じ手順の上に成り立つ
（違うのは形・既定色・角丸の付け方だけ）。この手順を_filled_shapeへ
集約し、それぞれは見た目の役割（コンテナ／装飾／罫線）だけを表す。

Connectorは2点を結ぶ線で、1点だけを装飾するMarkerとは本質的に別の
構造（つながりを表す）のため、_filled_shapeではなくadd_connectorの
上に別途実装する。
"""
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from .theme import LAYOUT, PALETTE


def _flat(shape, *, rounding=None, radius=None):
    shape.shadow.inherit = False
    effect_list = shape._element.spPr.find(qn("a:effectLst"))
    if effect_list is not None:
        shape._element.spPr.remove(effect_list)
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)
    if radius is not None and len(shape.adjustments):
        # ROUNDED_RECTANGLEのadj値は図形の短辺に対する割合として解釈されるため、
        # 同じ値でも縦横比が違う図形同士では角丸の見え方が揃わない。
        # 絶対の角丸半径を保つよう、短辺から都度adj値を逆算する。
        shape.adjustments[0] = min(radius / min(shape.width, shape.height), 0.5)
    elif rounding is not None and len(shape.adjustments):
        shape.adjustments[0] = rounding
    return shape


def _filled_shape(slide, shape_type, x, y, w, h, *, radius=None, rounding=None,
                  fill=None, line=None, line_width=Pt(0.7)):
    """図形を作り、既定のテーマスタイルを消し、塗り・枠線を設定する。

    Box/Marker/add_hairlineが共有する最小単位の手順。fill/line未指定は
    それぞれ透明/枠線なしになる。radius（絶対値）とrounding（相対値）は
    _flatにそのまま渡す。
    """
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    _flat(shape, radius=radius, rounding=rounding)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def Box(slide, x, y, w, h, *, rounded=True, radius=None, fill=None,
        line=None, line_width=Pt(0.7), elevated=False, shadow_fill=None):
    """矩形コンテナを描く（Atom層の基本図形）。

    rounded=Trueかつradius未指定ならradius_base（標準コンテナ用）を使う。
    画面の大部分を占める面にはradius=LAYOUT.radius_lgを明示する。
    rounded=Falseは角丸なしの区切り用の面（Background Zoneの既定）。

    elevated=Trueの場合、オフセットした面を背後に重ねて擬似的な影を
    作る。PowerPoint/LibreOffice/Keynote間で描画結果が揺れるネイティブ
    shadowエフェクト（outerShdw）は使わない（LibreOffice側で想定より
    濃く描画され、環境間で見え方が揺れることを確認済み）。fill/line
    を指定しない場合はそれぞれ塗りなし/枠線なしになる（区切りのため
    の透明な領域としても使える）。

    elevated=Trueで作った影と本体は、PowerPoint上でグループ化して
    1つの編集単位にする。グループ化しないと、手作業でサイズ変更する
    際に本体と影を毎回2つ別々に動かす必要があり、編集性が落ちるため。
    グループ化後もbox（本体）への参照はそのまま有効で、fill/line等を
    個別に変更できる。
    """
    if rounded:
        radius = LAYOUT.radius_base if radius is None else radius
    else:
        radius = None
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE

    if elevated:
        shadow = _filled_shape(slide, shape_type, x + Pt(1.5), y + Pt(2), w, h,
                               radius=radius, fill=shadow_fill or PALETTE.surface_subtle)

    box = _filled_shape(slide, shape_type, x, y, w, h, radius=radius,
                        fill=fill, line=line, line_width=line_width)

    if elevated:
        slide.shapes.add_group_shape([shadow, box])
    return box


def add_hairline(slide, x, y, w, *, color=PALETTE.grey_300, width=0.75):
    """細い横罫線（区切り線）を引く。テキストを持たない純粋な装飾図形。"""
    return _filled_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, Pt(width),
                         fill=color)


def _arrow_head(tag):
    # w/len（矢じりの幅・長さ）は明示指定する。省略時の既定値はレンダラー
    # 依存で、環境によって見え方がぶれるため（outerShdwと同じ理由）。
    element = OxmlElement(f"a:{tag}")
    element.set("type", "triangle")
    element.set("w", "lg")
    element.set("len", "lg")
    return element


def Connector(slide, x1, y1, x2, y2, *, style="straight", arrow="end",
             color=PALETTE.line_neutral, width=Pt(1.75)):
    """2点を結ぶ線（矢印可）を描く。

    横に並べたBox同士を「順番につながっている」と示すなど、要素間の
    関係を表す最小単位。Markerと違い2点間の関係を持つ点が異なる
    （Markerは1点の装飾、Connectorは2点をつなぐ線）。

    style: "straight"（直線、既定）または "elbow"（PowerPoint既定の
    自動経路で直角に曲がる線。始点・終点の相対位置から経路が決まる）。
    arrow: "end"（終点のみ矢印、既定）、"both"（両端）、"none"（矢印
    なし。区切り線ではなく関係を示す用途のまま矢印だけ外したい場合）。
    矢じりの形（triangle）・大きさ（lg/lg）はパラメータ化せず固定する。
    他の形状・大きさを必要とする実需が出てから検討する。
    """
    connector_type = MSO_CONNECTOR.ELBOW if style == "elbow" else MSO_CONNECTOR.STRAIGHT
    shape = slide.shapes.add_connector(connector_type, x1, y1, x2, y2)
    _flat(shape)
    shape.line.color.rgb = color
    shape.line.width = width
    if arrow != "none":
        line = shape._element.spPr.find(qn("a:ln"))
        if arrow == "both":
            line.append(_arrow_head("headEnd"))
        line.append(_arrow_head("tailEnd"))
    return shape


def Marker(slide, x, y, w, h, *, shape="bar", fill=PALETTE.line_brand, rounding=0.16):
    """テキストを持たない小さな単色アクセント図形（縦棒・点など）を描く。

    見出し脇の縦棒（SectionLead・章扉）、プロセスのゲート点など、
    「小さく・単色塗り・枠線なし」という同じ構造の装飾がrenderer側に
    個別実装されがちだったため、1箇所に集約する。Boxと違い枠線・影の
    バリエーションは持たない（装飾はBoxほど作り込まない、という区別）。

    shape="bar"（既定）は角丸の細い帯。shape="dot"は円。
    """
    shape_type = MSO_SHAPE.OVAL if shape == "dot" else MSO_SHAPE.ROUNDED_RECTANGLE
    return _filled_shape(slide, shape_type, x, y, w, h,
                         rounding=None if shape == "dot" else rounding, fill=fill)
