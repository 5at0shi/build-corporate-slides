"""スキルが標準搭載する部品を一覧できるショーケースデッキを生成する。

配色・タイポグラフィ・基本コンポーネント・アイコン・ネイティブチャート・
renderer 20種を、実際の描画関数を直接呼び出して本物の出力として並べる。
内容は常に実装と一致する（手書きの説明画像ではない）。デザインシステムや
renderer/icon/chartを追加・変更したら、このスクリプトを再実行して
参照を最新化する。

使い方（ワークスペース直下から）:
    ./.venv/bin/python .claude/skills/build-corporate-slides/scripts/build_capability_showcase.py
    ./.venv/bin/python .claude/skills/build-corporate-slides/scripts/render_and_check.py \\
        .claude/skills/build-corporate-slides/user-guide/capability-showcase.pptx \\
        --pdf .claude/skills/build-corporate-slides/user-guide/capability-showcase.pdf
"""
import sys
from pathlib import Path

SKILL = Path(__file__).resolve().parents[1]
WORKSPACE_ROOT = Path.cwd()
sys.path.insert(0, str(SKILL / "runtime" / "python"))

from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from slidekit import DeckBuilder, PALETTE  # noqa: E402
from slidekit.charts import add_native_chart  # noqa: E402
from slidekit.components import (add_background_zone, add_card,  # noqa: E402
                                 add_focus_panel, add_icon_list,
                                 add_item_list, add_key_message,
                                 add_section_lead)
from slidekit.icons import ICON_NAMES, add_icon  # noqa: E402
from slidekit.layout import Region  # noqa: E402
from slidekit.renderers import RENDERERS  # noqa: E402
from slidekit.theme import TYPE_BUSINESS, TYPE_DENSE, TYPE_LARGE_ROOM  # noqa: E402
from slidekit.typography import add_paragraph_textbox, add_text_list, add_textbox  # noqa: E402

FONT = "Hiragino Sans"


def caption(slide, x, y, w, text, *, size=10.5, bold=False, italic=False,
           color=PALETTE.text_secondary, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return box


def swatch(slide, x, y, w, h, color, name, sub):
    box = slide.shapes.add_shape(1, x, y, w, h)  # 1 = MSO_SHAPE.RECTANGLE
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    caption(slide, x, y + h + Inches(0.06), w, name, size=11, bold=True,
           color=PALETTE.text_primary)
    caption(slide, x, y + h + Inches(0.3), w, sub, size=9.5)


def centered_block(area, block_h, *, top_offset=Inches(0)):
    """area内、intro見出し等の下（top_offset）から、block_hの高さの帯を
    残りスペースの上下中央に配置したRegionを返す。

    項目数が少ないカタログページをtopのまま余白いっぱいへ伸ばすと、
    下部だけ大きく空いて間延びして見えるため（visual-quality.md）。
    """
    top = area.y + top_offset + max(0, (area.h - top_offset - block_h) // 2)
    return Region(area.x, top, area.w, block_h)


# ---------------------------------------------------------------- builder
builder = DeckBuilder.from_workspace(WORKSPACE_ROOT)

# ================================================================ 1. Cover
builder.add_cover(
    "build-corporate-slidesスキル 機能一覧",
    subtitle="配色・コンポーネント・アイコン・チャート・renderer 20種の実物見本",
    eyebrow="CAPABILITY SHOWCASE",
    brand_side="right", brand_shape="diagonal")

# ======================================================== 2. Color Palette
slide, area = builder.add_slide("配色（Semantic Color Roles）",
                                kicker="THEME.PY: PALETTE", page=2)
caption(slide, area.x, area.y, area.w,
       "色は役割（semantic role）で選ぶ。同じ役割は常に同じ色を使い、"
       "デッキ全体の一貫性を保つ。", size=12)

structural_heading_y = area.y + Inches(0.5)
caption(slide, area.x, structural_heading_y, area.w,
       "構造的な役割 — ブランドカラーに連動する既定の文字・線・面",
       size=11, bold=True, color=PALETTE.text_secondary)
grid_area = Region(area.x, structural_heading_y + Inches(0.34), area.w,
                   Inches(2.15))
grid = grid_area.rows([1, 1], gap=Inches(0.35))
colors = [
    ("text_primary", "見出し・本文の主文字"),
    ("text_secondary", "注記・補足"),
    ("line_neutral", "表・区切り・非強調の接続"),
    ("line_brand", "タイトル下線・セクションマーカー"),
    ("surface_subtle", "ニュートラルな補助面"),
    ("surface_brand_soft", "青系のまとまり・選択範囲"),
    ("surface_teal_soft", "異なる構造の補助系列"),
    ("focus_primary", "最重要項目"),
    ("accent_secondary", "構造を区別する限定的アクセント"),
]
for row_index, row in enumerate(grid):
    cols = row.columns([1] * 5, gap="standard")
    for col_index, region in enumerate(cols):
        index = row_index * 5 + col_index
        if index >= len(colors):
            break
        name, desc = colors[index]
        swatch(slide, region.x, region.y, region.w, Inches(0.7),
              getattr(PALETTE, name), name, desc)

sign_heading_y = grid_area.y + grid_area.h + Inches(0.4)
caption(slide, area.x, sign_heading_y, area.w,
       "符号・重大度ロール — ブランドカラーとは独立。良い/悪い/注意の"
       "意味でのみ使い、装飾としては使わない（背景の淡色版はp.4）",
       size=11, bold=True, color=PALETTE.text_secondary)
sign_row = Region(area.x, sign_heading_y + Inches(0.34), area.w, Inches(1.06))
sign_cols = sign_row.columns([1] * 5, gap="standard")
sign_colors = [
    ("positive", "良い結果・達成・承認"),
    ("negative", "悪い結果・リスク・未達"),
    ("warning", "注意・要確認"),
]
for offset, (name, desc) in enumerate(sign_colors, start=1):
    region = sign_cols[offset]
    swatch(slide, region.x, region.y, region.w, Inches(0.7),
          getattr(PALETTE, name), name, desc)

# ======================================================= 3. Typography
slide, area = builder.add_slide("タイポグラフィ（3モード）",
                                kicker="THEME.PY: typography_for()", page=3)
caption(slide, area.x, area.y, area.w,
       "deck.modeで選ぶ。businessの本文は12ptを基準にする。"
       "denseは情報量が多いページだけに使い、large-roomは遠距離投影が"
       "明示された場合だけ使う（deck.mode自体はbusinessのまま個別ページに"
       "density: denseを指定できる）。以下は実際のpt数で描いた実寸サンプル。",
       size=11.5)
modes = [
    ("business（標準）", TYPE_BUSINESS,
     "個人PC閲覧・事前配布・画面共有が前提。通常の社内資料はすべてこれ。"),
    ("dense", TYPE_DENSE,
     "表・比較など情報量が多いページだけに使う。資料全体はdenseにしない。"),
    ("large-room", TYPE_LARGE_ROOM,
     "大会議室・講演など遠距離投影が明示された場合だけ使う。"),
]
mode_block = Region(area.x, area.y + Inches(0.86), area.w, Inches(4.5))
mode_cols = mode_block.columns([1, 1, 1], gap="wide")
sample_rows = [
    ("Title", "サンプル見出し", "title"),
    ("Section", "セクション見出し", "section"),
    ("Body", "本文サンプル Sample 123", "body"),
    ("Small", "注記・出典サンプル", "small"),
]
for (label, type_, when), region in zip(modes, mode_cols):
    add_background_zone(slide, region.x, region.y, region.w, mode_block.h,
                        tone="neutral", rounded=True)
    inner = region.inset(Inches(0.26), Inches(0.22))
    y = inner.y
    add_textbox(slide, inner.x, y, inner.w, Inches(0.3), label,
               size=Pt(14), color=PALETTE.text_primary, bold=True, font=FONT)
    y += Inches(0.36)
    add_paragraph_textbox(slide, inner.x, y, inner.w, Inches(0.6), [
        {"segments": [(when, {"size": Pt(9.5), "color": PALETTE.text_secondary,
                              "font": FONT})], "line_spacing": 1.2}])
    y += Inches(0.62)
    for row_label, sample_text, attr in sample_rows:
        size = getattr(type_, attr)
        add_textbox(slide, inner.x, y, inner.w, Inches(0.16),
                   f"{row_label}  {size.pt:g}pt", size=Pt(8.5),
                   color=PALETTE.blue, bold=True, font=FONT)
        y += Inches(0.2)
        sample_h = Inches(size.pt / 72 + 0.08)
        add_textbox(slide, inner.x, y, inner.w, sample_h,
                   sample_text, size=size,
                   color=PALETTE.text_primary,
                   bold=(attr in ("title", "section")), font=FONT)
        y += sample_h + Inches(0.16)

# =============================================== 4. Components: 面・パネル
slide, area = builder.add_slide("コンポーネント（面・パネル）",
                                kicker="COMPONENTS.PY", page=4)
block4 = centered_block(area, Inches(4.0), top_offset=Inches(0.15))
row = block4.rows([2.0, 1.5], gap=Inches(0.5))
top_cols = row[0].columns([1, 1, 1], gap="standard")
add_card(slide, top_cols[0].x, top_cols[0].y, top_cols[0].w, Inches(1.5))
caption(slide, top_cols[0].x, top_cols[0].y + Inches(1.58), top_cols[0].w,
       "add_card()", size=11, bold=True, color=PALETTE.text_primary)
caption(slide, top_cols[0].x, top_cols[0].y + Inches(1.8), top_cols[0].w,
       "独立して比較・操作する単位。細線または軽い影。", size=9.5)

add_background_zone(slide, top_cols[1].x, top_cols[1].y, top_cols[1].w,
                    Inches(1.5), tone="brand-soft", rounded=True)
caption(slide, top_cols[1].x, top_cols[1].y + Inches(1.58), top_cols[1].w,
       'add_background_zone(tone="brand-soft")', size=11, bold=True,
       color=PALETTE.text_primary)
caption(slide, top_cols[1].x, top_cols[1].y + Inches(1.8), top_cols[1].w,
       "緩やかなまとまり・範囲・フェーズ。境界・影なし。", size=9.5)

add_focus_panel(slide, top_cols[2].x, top_cols[2].y, top_cols[2].w, Inches(1.5),
                tone="solid")
caption(slide, top_cols[2].x, top_cols[2].y + Inches(1.58), top_cols[2].w,
       'add_focus_panel(tone="solid")', size=11, bold=True,
       color=PALETTE.text_primary)
caption(slide, top_cols[2].x, top_cols[2].y + Inches(1.8), top_cols[2].w,
       "ページ内の重点対象。色・線・サイズで明確化。", size=9.5)

caption(slide, row[1].x, row[1].y, row[1].w,
       'add_background_zone()のtone一覧（符号・重大度の淡色も汎用のComponent部品として使える）',
       size=10, bold=True, color=PALETTE.text_secondary)
bottom_cols = row[1].inset(top=Inches(0.3)).columns([1, 1, 1, 1, 1], gap="tight")
tones = ["neutral", "teal-soft", "positive-soft", "negative-soft", "warning-soft"]
for i, tone in enumerate(tones):
    region = bottom_cols[i]
    add_background_zone(slide, region.x, region.y, region.w, Inches(1.0),
                        tone=tone, rounded=True)
    caption(slide, region.x, region.y + Inches(1.08), region.w,
           f'tone="{tone}"', size=10, bold=True, color=PALETTE.text_primary)

# ============================================ 5. Components: テキスト・リスト
slide, area = builder.add_slide("コンポーネント（テキスト・リスト）",
                                kicker="COMPONENTS.PY", page=5)
caption(slide, area.x, area.y, area.w,
       "同じ意味のまとまりは、行ごとに別shapeへ分割せず一つの複数段落"
       "textboxにまとめる。選択・移動・順序の入れ替えを一つの操作で"
       "行えるようにするための編集性への配慮。", size=12)
block5 = centered_block(area, Inches(3.75), top_offset=Inches(0.55))
row5 = block5.rows([2.5, 1.1], gap=Inches(0.4))
top = row5[0]
lead_col, list_col, icon_col = top.columns([0.8, 1, 1], gap="wide")
add_section_lead(slide, lead_col.x, lead_col.y, lead_col.w, "add_section_lead()")
add_item_list(slide, lead_col.x, lead_col.y + Inches(0.6), lead_col.w, Inches(1.6), [
    {"title": "見出し + 項目", "body": "add_item_listは複数項目を一つのtextboxへまとめる"}])
caption(slide, lead_col.x, lead_col.y + Inches(2.3), lead_col.w,
       "add_section_lead() + add_item_list()", size=10, bold=True,
       color=PALETTE.text_primary)

add_text_list(slide, list_col.x, list_col.y, list_col.w, Inches(1.9), [
    {"title": "marker=\"number\"", "body": "01, 02...を自動採番する"},
    {"title": "divider=True", "body": "行ごとに安全な位置へ罫線を引く"}],
    marker="number", divider=True, max_row_h=Inches(0.9))
caption(slide, list_col.x, list_col.y + Inches(1.9), list_col.w,
       'add_text_list(marker="number", divider=True)', size=10, bold=True,
       color=PALETTE.text_primary)

add_icon_list(slide, icon_col.x, icon_col.y, icon_col.w, Inches(1.7), [
    "add_icon_listはアイコンを行ごとに独立させる", "文章は一つのtextboxへまとめる"],
    icon="check")
caption(slide, icon_col.x, icon_col.y + Inches(1.9), icon_col.w,
       "add_icon_list()", size=10, bold=True, color=PALETTE.text_primary)

msg_row = row5[1]
msg_cols = msg_row.columns([1, 1, 1, 1], gap="standard")
for (style, label), region in zip(
        [("editorial", "editorial"), ("subtle", "subtle"),
         ("solid", "solid"), ("card", "card")], msg_cols):
    h = Inches(0.82)
    add_key_message(slide, region.x, region.y, region.w,
                    f"Key Message: {label}", style=style)
    caption(slide, region.x, region.y + h + Inches(0.08), region.w,
           f'add_key_message(style="{style}")', size=9.5)

# ==================================================================== 6. Icons
slide, area = builder.add_slide("アイコン（14種）", kicker="ICONS.PY: add_icon()",
                                page=6)
caption(slide, area.x, area.y, area.w,
       "外部アセットを使わず、python-pptxの図形だけで描くベクターアイコン。"
       "状態・分類を一目で区別させたい場合だけ使う。", size=12)
names = sorted(ICON_NAMES)
icon_block = centered_block(area, Inches(3.0), top_offset=Inches(0.55))
icon_grid = icon_block.rows([1, 1], gap=Inches(0.5))
for row_index, row in enumerate(icon_grid):
    cols = row.columns([1] * 7, gap="tight")
    for col_index, region in enumerate(cols):
        index = row_index * 7 + col_index
        if index >= len(names):
            break
        name = names[index]
        cx = region.x + region.w // 2 - Inches(0.4)
        add_icon(slide, cx, region.y, Inches(0.8), name, color=PALETTE.blue)
        caption(slide, region.x, region.y + Inches(0.95), region.w, name,
               size=10.5, align=PP_ALIGN.CENTER)

# ============================================================ 7. Charts
slide, area = builder.add_slide("ネイティブチャート（6種）",
                                kicker="CHARTS.PY: add_native_chart()", page=7)
caption(slide, area.x, area.y, area.w,
       "PNG画像と異なり、貼り付け後もPowerPoint上で数値・系列名を直接編集"
       "できる。chart_with_insightのchartフィールドから使う。", size=12)
chart_grid = area.inset(top=Inches(0.55)).rows([1, 1], gap=Inches(0.35))
chart_specs = [
    ("column", {"chart_type": "column", "categories": ["Q1", "Q2", "Q3"],
               "series": [{"name": "s1", "values": [4, 7, 6]}]}),
    ("stacked_column", {"chart_type": "stacked_column",
                        "categories": ["Q1", "Q2", "Q3"],
                        "series": [{"name": "既存", "values": [4, 5, 6]},
                                   {"name": "新規", "values": [2, 3, 4]}]}),
    ("bar", {"chart_type": "bar", "categories": ["A", "B", "C"],
            "series": [{"name": "s1", "values": [5, 8, 3]}]}),
    ("line", {"chart_type": "line", "categories": ["4月", "5月", "6月"],
             "series": [{"name": "s1", "values": [10, 14, 12]}]}),
    ("pie", {"chart_type": "pie", "categories": ["A", "B", "C"],
            "series": [{"name": "s1", "values": [5, 3, 2]}]}),
    ("scatter", {"chart_type": "scatter", "series": [
        {"name": "s1", "points": [{"x": 1, "y": 2}, {"x": 2, "y": 3.4},
                                  {"x": 3, "y": 1.8}]}]}),
]
for row_index, row in enumerate(chart_grid):
    cols = row.columns([1, 1, 1], gap="standard")
    for col_index, region in enumerate(cols):
        index = row_index * 3 + col_index
        if index >= len(chart_specs):
            break
        label, kwargs = chart_specs[index]
        chart_h = region.h - Inches(0.3)
        add_native_chart(slide, region.x, region.y, region.w, chart_h,
                         typography=TYPE_BUSINESS, **kwargs)
        caption(slide, region.x, region.y + chart_h + Inches(0.04), region.w,
               f'chart.type: "{label}"', size=10, bold=True,
               color=PALETTE.text_primary, align=PP_ALIGN.CENTER)


# =================================================== 8-26. Renderer Catalog
# renderer-catalog.mdの表と同じ順序で並べ、ドキュメントと実物の並び順を
# 揃える。variant違いは各typeの基本例の直後に置く（末尾にまとめない）。
# 「使う状況」列と同じ文言を使い、1ページのcontent_region()内に収まる、
# 控えめな現実的な分量で作る（stress-test用の極端な分量ではない）。
FOOTER_Y = Inches(7.5) - Inches(0.32)


def renderer_footer(type_name, desc):
    slide = builder.presentation.slides[-1]
    caption(slide, Inches(0.67), FOOTER_Y, Inches(11.99),
           f"renderer: “{type_name}” — {desc}", size=10,
           italic=True)


# (renderer_type, footerラベル, footer説明, spec) を表示順に並べる。
catalog = [
    ("comparison", "comparison", "二つの観点、価値とリスク、現状と将来", {
        "title": "comparison: 二つの観点を公平に比較する", "density": "standard",
        "primary_message": "balancedは対等な2列、asymmetricは片側を補助情報として弱める",
        "left": {"heading": "向いている場面", "items": [
            {"title": "価値とリスク", "body": "メリットとデメリットを両方読ませたいとき"},
            {"title": "現状と将来", "body": "Before/Afterを対等な分量で見せたいとき"}]},
        "right": {"heading": "編集の単位", "items": [
            {"title": "列ごとに1つのtextbox", "body": "項目を行ごとに別shapeへ分割しない"},
            {"title": "balanced / asymmetric", "body": "片側を弱めたい場合はasymmetricを選ぶ"}]},
    }),
    ("evidence_and_decision", "evidence_and_decision", "根拠から一つの推奨判断を導く", {
        "title": "evidence_and_decision: 根拠から一つの結論を導く", "density": "standard",
        "primary_message": "結論は1文に絞り、根拠は独立した領域で支える",
        "evidence_heading": "使う場面", "evidence": [
            {"title": "根拠が複数、結論は1つ", "body": "複数の事実から単一の推奨判断を示したいとき"},
            {"title": "判断理由を明示したい", "body": "結論だけでなく根拠も同じページで示したいとき"}],
        "decision_heading": "推奨する使い方",
        "decision_detail": "evidenceは箇条書きの根拠群、decisionは結論文。結論を薄めないよう1文に絞る。",
    }),
    ("scope_and_exclusions", "scope_and_exclusions", "対象範囲と対象外を同時に示す", {
        "title": "scope_and_exclusions: 対象と対象外を同時に示す", "density": "dense",
        "primary_message": "対象範囲と除外範囲を1ページで誤解なく伝える",
        "scope_heading": "対象に含めるもの（scope）", "scope": [
            {"label": "IN", "title": "scope配列", "body": "対象とする意味単位を列挙する"},
            {"label": "OUT", "title": "exclusionsは別枠", "body": "対象外と混在させず一括リストにする"}],
        "period": "period: 期間や数値などの補足情報を1行で添えられる",
        "exclusions_heading": "対象外（exclusions）",
        "exclusions": ["除外項目は配列でまとめて一括表示する", "個別のcard化はしない",
                       "対象と対象外は必ず両方示す"],
    }),
    ("process_with_gates", "process_with_gates", "フェーズ、作業、判断時点", {
        "title": "process_with_gates: フェーズと承認ゲートを分けて描く", "density": "dense",
        "primary_message": "フェーズは横並びのレーン、ゲートは別レイヤーの時系列軸",
        "phases": [
            {"title": "フェーズ", "label": "PHASE", "weight": 1,
             "items": ["phases配列の各要素が1つのレーンになる", "weightで幅の比率を調整できる"]},
            {"title": "ゲート", "label": "GATE", "weight": 1,
             "items": ["gatesは0〜1のpositionで時系列上に配置", "承認・判断のタイミングを示す"]}],
        "gates": [{"title": "ゲート1", "position": 0.3},
                  {"title": "ゲート2", "position": 1.0}],
    }),
    ("table_with_conclusion", "table_with_conclusion",
     "条件、比較、評価を表で読む。ネイティブtableのため件数が増えても行を足すだけ", {
        "title": "table_with_conclusion: 表と結論をセットで示す", "density": "dense",
        "primary_message": "結論は表の外に独立させ、表そのものは常にネイティブPowerPoint tableで作る",
        "columns": [{"key": "item", "label": "項目", "weight": 1.3},
                    {"key": "desc", "label": "編集性の工夫", "weight": 3.3},
                    {"key": "note", "label": "備考", "weight": 0.9, "align": "center"}],
        "rows": [
            {"item": "ネイティブtable", "desc": "PowerPointの表機能で作るため、画像化せずセルの文字を直接クリックして編集できる", "note": "必須"},
            {"item": "縦線なし", "desc": "列の区切りは余白だけで示す。全辺に線を引くと構造がのっぺりして見える", "note": "任意"},
            {"item": "ヘッダー下だけ強罫線", "desc": "本文行は薄い横線のみ。ヘッダーと本文の境界だけを強く見せる", "note": "任意"},
            {"item": "_highlight: true", "desc": "行に付けると太字＋淡色背景で強調表示。装飾用の追加shapeは不要",
             "note": "任意", "_highlight": True},
            {"item": "rows配列", "desc": "各要素が1行になる。件数の増減はrows配列を増減するだけでレイアウト再設計は不要", "note": "必須"}],
    }),
    ("table_with_insight", "table_with_insight",
     "表から複数の気づきを箇条書きで示す（chart_with_insightの表版）", {
        "title": "table_with_insight: 四半期指標の推移から複数の気づきを示す", "density": "dense",
        "primary_message": "Q3は量を優先した結果、質が犠牲になっている可能性がある。次四半期は受注率の回復を優先指標に据える",
        "insight_heading": "読み取れること",
        "insights": [
            "Q2は商談数が減った一方、受注率は最も高かった（質を優先した営業活動の効果）",
            "Q3は商談数が増えたが受注率が低下しており、対応の丁寧さが薄れている可能性がある",
            "平均商談期間はQ2に短縮したが、Q3で再び伸びている"],
        "columns": [{"key": "metric", "label": "指標", "weight": 1.4},
                    {"key": "q1", "label": "Q1", "weight": 1},
                    {"key": "q2", "label": "Q2", "weight": 1},
                    {"key": "q3", "label": "Q3", "weight": 1}],
        "rows": [
            {"metric": "新規商談数", "q1": "120件", "q2": "98件", "q3": "145件"},
            {"metric": "平均商談期間", "q1": "32日", "q2": "29日", "q3": "35日"},
            {"metric": "受注率", "q1": "18%", "q2": "22%", "q3": "15%", "_highlight": True}],
    }),
    ("chart_with_insight", "chart_with_insight (standard)", "グラフ全体を公平に読む", {
        "title": "chart_with_insight (standard): グラフと読み取りを公平に示す",
        "density": "standard", "variant": "standard",
        "primary_message": "グラフ全体を公平に読ませたいときはstandardを選ぶ",
        "insight_heading": "使う場面",
        "insights": ["グラフの全体像を偏りなく見せたいとき",
                     "column/stacked_column/bar/line/pie/scatterに対応",
                     "chartを指定すると数値をPowerPoint上で直接編集できる"],
        "chart": {"type": "column", "categories": ["項目1", "項目2", "項目3"],
                 "series": [{"name": "系列A", "values": [12, 38, 42]}]},
    }),
    ("chart_with_insight", "chart_with_insight (conclusion-led)", "一つの主張を強く伝える", {
        "title": "chart_with_insight (conclusion-led): 一つの主張を強く伝える",
        "density": "standard", "variant": "conclusion-led",
        "primary_message": "1つの主張を強く伝えたいときはconclusion-ledを選ぶ",
        "insight_heading": "使う場面",
        "insights": ["グラフから読み取れる主張が1つに絞れるとき",
                     "primary_messageで結論を先に示し、グラフで裏付ける"],
        "chart": {"type": "line", "categories": ["項目1", "項目2", "項目3", "項目4"],
                 "series": [{"name": "系列A", "values": [5, 22, 38, 42]}]},
    }),
    ("org_layers", "org_layers", "意思決定・運営など縦の責任階層と、横に並ぶ実行部門", {
        "title": "org_layers: 縦の階層と横の実行部門を分けて描く", "density": "standard",
        "primary_message": "layersは2件までを目安にする（本文がはみ出さないように）",
        "layers": [
            {"heading": "階層1", "title": "layers配列", "body": "上から順に責任階層のバンドとして積む"},
            {"heading": "階層2", "title": "2件が上限の目安", "body": "3件以上は各層の本文が収まらずpreflightが警告する"}],
        "execution_heading": "実行部門（execution）",
        "execution": [{"title": "横に並ぶCard", "body": "execution配列の各要素が1枚のCardになる"},
                      {"title": "部門数に応じ自動分割", "body": "列数は要素数に合わせて均等割りする"}],
    }),
    ("priority_actions", "priority_actions", "優先度付きの課題と、対応する方針", {
        "title": "priority_actions: 優先度付きの課題と対応策を並べる", "density": "standard",
        "primary_message": "issuesは優先度ラベル付き、actionsは対応方針の一括リスト",
        "issues": [
            {"priority": "最優先", "title": "issues配列", "body": "priority/title/bodyを持つ課題を優先度順に並べる"},
            {"priority": "中", "title": "優先度ラベル", "body": "最優先・高・中など任意の文字列を指定できる"}],
        "actions": ["actionsは対応方針の一括リスト", "淡色パネルの中に箇条書きで表示する",
                    "issuesとは別領域で独立して読める"],
    }),
    ("stage_track", "stage_track",
     "段階的な進行を同格のCardで示す。各段はbody（文章）でもitems（箇条書き）でも書ける（期間の幅を示すならtimeline）", {
        "title": "stage_track: 段階的な進行を同格のCardで示す", "density": "standard",
        "primary_message": "stagesは配列の順に横並びの同格Cardになる。中身はbody（文章）とitems（箇条書き）のどちらでも書ける",
        "connectors": False,
        "stages": [
            {"label": "STEP1", "title": "stages配列", "body": "label/title/bodyを持つ段階を順に並べる"},
            {"label": "STEP2", "title": "同格のCard", "body": "各段階は優劣なく並列に見せる"},
            {"label": "STEP3", "title": "itemsで箇条書き",
             "items": ["段ごとにやることを並べる", "bodyと併用もできる"]}],
    }),
    ("stage_track", "stage_track (connectors: true)",
     "矢印表現はここだけの限定的なオプション。cycleの円周矢印は構造上必須のため別枠", {
        "title": "stage_track (connectors: true): 矢印で順序を明示する",
        "density": "standard",
        "connectors": True,
        "primary_message": "段階の間に順序・因果があることを強調したい場合だけconnectors:trueにする",
        "stages": [
            {"label": "STEP1", "title": "既定はconnectors: false", "body": "他の型の基本例と同じ矢印なしの見た目"},
            {"label": "STEP2", "title": "矢印はここでだけ使う", "body": "この構造図に限り、順序性を強めたいときに選べるオプション"},
            {"label": "STEP3", "title": "太さ・矢じりは固定", "body": "Atom層のConnectorが一貫した太さで描く"}],
    }),
    ("timeline", "timeline",
     "開始と終了が異なる取り組みが並走する計画（ロードマップ、ガント）。帯の幅が期間を表す（stage_trackは順序、process_with_gatesは判断時点）", {
        "title": "timeline: 期間ぶんの幅を持つ帯を共通の時間軸へ並べる", "density": "standard",
        "primary_message": "periodsが時間軸の目盛り、rowsのstart/endが占める区間（1始まり・endを含む）。範囲外の値は軸の端へ丸める",
        "periods": ["Q1", "Q2", "Q3", "Q4"],
        "rows": [
            {"label": "rows配列", "title": "1行が1つの取り組みになる", "start": 1, "end": 2},
            {"label": "start / end", "title": "占める区間を期間番号で指定する", "start": 2, "end": 3,
             "tone": "teal"},
            {"label": "tone", "title": "帯の色は意味で選ぶ", "start": 3, "end": 4, "tone": "warning"},
            {"label": "自動回り込み",
             "title": "帯の中に収まらない長さのtitleは、文字が潰れないよう自動的に帯の右側へ回る",
             "start": 1, "end": 1, "tone": "neutral"}],
    }),
    ("numbered_list", "numbered_list", "アジェンダ、依頼事項など番号付きの単列項目", {
        "title": "numbered_list: 番号付きの単列項目を並べる", "density": "standard",
        "primary_message": "message_positionでtop/bottomの導入文・結論の位置を選べる",
        "message_position": "bottom",
        "items": [
            {"title": "items配列", "body": "番号は自動で振られ、titleとbodyを持てる"},
            {"title": "アジェンダやNext Stepに最適", "body": "単列の番号付き項目だけで構成されるページに使う"},
            {"title": "row_hは項目数で自動調整", "body": "項目が少ないほど各行の余白が広がる"}],
    }),
    ("matrix", "matrix", "2軸で選択肢を4象限に整理する（ポートフォリオ分析等）", {
        "title": "matrix: 2軸で4象限に整理する", "density": "standard",
        "primary_message": "軸上の正確な数値でなく、象限のどれに属するかが要点のときに使う",
        "x_axis": {"low": "x_axis.low", "high": "x_axis.high"},
        "y_axis": {"low": "y_axis.low", "high": "y_axis.high"},
        "cells": [
            {"label": "強調", "title": "cells配列", "body": "rows×cols件ちょうど必要。各マスは独立編集可能", "emphasis": True},
            {"label": "補足", "title": "emphasis:Trueで強調", "body": "1マスだけ強調したい場合に使う"},
            {"label": "補足", "title": "正確な座標が要点なら", "body": "chart_with_insightのscatterを使う"},
            {"label": "補足", "title": "背景はBackground Zone", "body": "マスごとに独立した面として描く"}],
    }),
    ("matrix", "matrix (x_axis/y_axis省略)",
     "SWOT等、軸のない固定カテゴリの整理。rows/colsを指定すれば2x2を超えるマトリクスも同じ構造で作れる", {
        "title": "matrix (x_axis/y_axis省略): 軸のない4カテゴリで示す", "density": "standard",
        "primary_message": "SWOT等、連続軸を持たない固定4カテゴリはx_axis/y_axisを省略すると軸ラベル分の余白を返す",
        "cells": [
            {"label": "S", "title": "強み", "body": "x_axis/y_axisを省くと軸キャプション行を省く"},
            {"label": "W", "title": "弱み", "body": "cellsの構造は軸ありと同じ"},
            {"label": "O", "title": "機会", "body": "SWOT等、固定4カテゴリの整理に使う"},
            {"label": "T", "title": "脅威", "body": "連続軸上の位置づけが要点ならx_axis/y_axisを指定"}],
    }),
    ("issue_tree", "issue_tree",
     "論点をMECEに分解する（イシューツリー、ロジックツリー）。線に矢印は付けない（流れではなく包含関係のため）", {
        "title": "issue_tree: 論点を3階層に分解して示す", "density": "standard",
        "primary_message": "rootが分解する問い、branchesが第1階層、各branchのitemsが内訳。階層は3段固定で、深い分解はページを分ける",
        "root": {"label": "root", "title": "分解する問いを置く",
                "body": "label / title / bodyの3項目"},
        "branches": [
            {"title": "branches配列", "body": "第1階層の分解軸。5件までが目安",
             "items": ["itemsが内訳になる", "1件が1行"]},
            {"title": "itemsは合計12件まで", "body": "超えると1行の高さが狭くなり警告する",
             "items": ["枝ごとに件数が違ってよい"]},
            {"title": "items省略時", "body": "どの枝もitemsを持たない場合は2段のツリーになり、枝が広くなる",
             "items": ["この例では3段"]}],
    }),
    ("stat_highlight", "stat_highlight", "単一の実績数値を主役に、補足指標と結論を示す", {
        "title": "stat_highlight: 単一の実績数値を主役にする", "density": "standard",
        "primary_message": "stat.valueが主指標、supportingは補足指標のCard群",
        "stat": {"value": "1", "label": "主役にする数値は1つだけ",
                "detail": "stat.value / label / detailの3項目で構成する"},
        "supporting": [{"value": "N件", "label": "supportingは補足指標"},
                       {"value": "任意", "label": "複数観点の対等比較はtable_with_conclusionを優先"}],
    }),
    ("stat_highlight", "stat_highlight (KPI dashboard)",
     "toneは数値の符号でなく意味で指定する（-42%でもtone: positiveになり得る）", {
        "title": "stat_highlight (stat省略): KPIダッシュボードとして一覧する",
        "density": "standard",
        "primary_message": "statを省略すると、supportingが均等グリッドのKPIダッシュボードになる",
        "supporting": [
            {"value": "94%", "label": "満足度"},
            {"value": "-42%", "label": "作業時間削減率（削減=良い結果）", "tone": "positive"},
            {"value": "+8%", "label": "解約率（増加=悪い結果）", "tone": "negative"},
            {"value": "12", "label": "導入部署数"}],
    }),
    ("waterfall", "waterfall",
     "AからBへの変化を増減要因に分解する（ブリッジ図、EBITDAウォーク、予算差異、価格/数量/構成の分解）", {
        "title": "waterfall: 変化した理由を増減の棒へ分解する", "density": "standard",
        "primary_message": "kind: totalの棒は基準線0からの絶対値、kindを省略した棒は直前までの累計に対する増減として宙に浮く",
        "bars": [
            {"label": "開始値（kind: total）", "value": 100, "kind": "total",
             "value_label": "100"},
            {"label": "増加要因（符号で自動着色）", "value": 34, "value_label": "+34"},
            {"label": "減少要因", "value": -18, "value_label": "-18"},
            {"label": "減だがtone: positive（削減=良い結果）", "value": -9,
             "value_label": "-9", "tone": "positive"},
            {"label": "終了値（kind: total）", "value": 107, "kind": "total",
             "value_label": "107"}],
    }),
    ("funnel", "funnel", "順を追って絞り込まれていく推移（営業パイプライン、市場規模のTAM/SAM/SOM等）。insights指定で使用ケースや気づきを併記できる", {
        "title": "funnel: 絞り込みの推移を帯の幅で示す",
        "density": "standard",
        "primary_message": "stagesは値の大きい順に並べる。帯の幅はおおよその絞り込み具合を示す構造表現で、正確な比率が要点ならchart_with_insightの棒グラフを使う",
        "insight_heading": "使う場面",
        "insights": [
            "リード獲得の営業パイプライン（問い合わせ→商談→提案→受注等、下図は例）",
            "市場規模の絞り込み（TAM→SAM→SOM）",
            "帯の幅は値に比例するのではなく、平方根で圧縮したおおよその縮み具合を表す",
            "insights省略時はこの箇条書きをやめ、帯が全幅を使う標準の見た目になる"],
        "stages": [
            {"title": "問い合わせ", "value": 420, "value_label": "420件"},
            {"title": "商談", "value": 95, "value_label": "95件"},
            {"title": "提案", "value": 38, "value_label": "38件"},
            {"title": "受注", "value": 12, "value_label": "12件"}],
    }),
    ("cycle", "cycle", "繰り返し・循環するプロセス（PDCA等）", {
        "title": "cycle: 繰り返すプロセスを円周上のCard群で示す", "density": "standard",
        "primary_message": "stepsを時計回りに配置し、最後尾から先頭へも矢印で結んで輪にする",
        "steps": [
            {"label": "STEP1", "title": "steps配列"},
            {"label": "STEP2", "title": "円周上に配置"},
            {"label": "STEP3", "title": "隣接を矢印で連結"},
            {"label": "STEP4", "title": "4〜6件が目安"}],
    }),
]

page = 8
RENDERERS["section_divider"](builder, {
    "title": "renderer カタログ（20種）", "kicker": "SECTION DIVIDER",
    "subtitle": "各typeの最小構成の例。件数・文言・配色は指示に応じて自由に調整でき、"
               "これが唯一の形ではない",
}, page)
renderer_footer("section_divider", "複数テーマを扱う資料の章区切り")
page += 1

for renderer_type, label, desc, spec in catalog:
    RENDERERS[renderer_type](builder, spec, page)
    renderer_footer(label, desc)
    page += 1

output = builder.save(SKILL / "user-guide" / "capability-showcase.pptx")
print("PPTX:", output)
