"""ページ内に置く部品のうち、ユーザー・AIが名前で指示する対象。
Atom層に属する（層構成の全体像はARCHITECTURE.mdを参照）。

「コンポーネント」は層の名前ではなく、Atom層のうち公開している部品の
呼び名。capability-showcaseに実物を載せ、「4ページ目のBackground Zone
を広げて」のように名指しできるのはここに定義された名前だけ。

中身は2種類ある。どちらもAtom層であり、両者の間に階層はない。
- Atomの名前付きプリセット: add_card / add_background_zone /
  add_focus_panel はいずれもatoms.pyのBoxへ引数を変えて委譲するだけ。
- 少数のAtomの定型的な組み合わせ: add_section_lead（Marker＋テキスト）、
  add_key_message（罫線またはBox＋テキスト）など。

add_item_list / add_icon_listは実体をtypography.pyのadd_text_listへ
委譲する（テキスト処理をtypography.pyへ集約しているため）。

表紙・章扉・ページヘッダーはページ内の部品ではなくページ外枠のため、
pageframe.pyに置く。
"""
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from .atoms import Box, Marker, add_hairline
from .layout import Region
from .theme import LAYOUT, PALETTE
from .typography import (_type_for, add_text_list, add_textbox, set_run,
                         style_text_frame)


# add_section_leadの縦棒マーカー分として、呼び出し側が下に置く要素へ
# 空けるべき最小オフセット。マーカー高さが呼び出し側でも変わる場合は
# marker_h + SECTION_LEAD_GAPを使い、値がズレて重なるのを防ぐ。
SECTION_LEAD_GAP = Inches(0.1)


def add_section_lead(slide, x, y, w, text, *, color=PALETTE.line_brand,
                     size=None, marker_h=Inches(0.38)):
    typography = _type_for(slide)
    size = size or typography.section
    marker = Marker(slide, x, y, Inches(0.06), marker_h, fill=color)
    label = add_textbox(slide, x + Inches(0.16), y - Inches(0.01),
                        w - Inches(0.16), Inches(0.42), text,
                        size=size, color=PALETTE.text_primary, bold=True)
    return marker, label


def add_background_zone(slide, x, y, w, h, *, tone="brand-soft",
                        rounded=False):
    tones = {
        "brand-soft": PALETTE.surface_brand_soft,
        "neutral": PALETTE.surface_subtle,
        "teal-soft": PALETTE.surface_teal_soft,
        "positive-soft": PALETTE.surface_positive_soft,
        "negative-soft": PALETTE.surface_negative_soft,
        "warning-soft": PALETTE.surface_warning_soft,
    }
    if tone not in tones:
        raise ValueError(f"未定義のtoneです: {tone}")
    # 画面の大部分を占める面のため、角丸を使う場合はradius_lg（Card等の
    # radius_baseより大きい半径）を使う。小さい半径だと丸みが足りずに見える。
    return Box(slide, x, y, w, h, rounded=rounded, radius=LAYOUT.radius_lg,
               fill=tones[tone], line=None)


def add_emphasis_zone(slide, region, *, tone="neutral"):
    """columns()等で切り出した列を、外側へわずかにはみ出す淡色の面で
    包み、隣接する列より強調して見せる（comparisonのasymmetric variant
    とpriority_actionsの対応方針列が共有していた同一の操作を集約した）。

    上下は非対称（上0.08in・下0.04in）にはみ出す。section_leadの見出し
    マーカーの上に余白を作りつつ、下は詰めすぎないための調整。
    """
    outer = region.inset(Inches(-0.18), Inches(-0.08), Inches(-0.18), Inches(-0.04))
    return add_background_zone(slide, outer.x, outer.y, outer.w, outer.h,
                               tone=tone, rounded=True)


def add_panel(slide, x, y, w, h, *, tone="neutral", rounded=True,
             inset_x=Inches(0.32), inset_y=Inches(0.24)):
    """背景面を描き、内側の余白を差し引いたRegionを返す。

    見出し+リストや複数段落のブロックを「意図的な余白を持つ面」として
    見せたいときに使う。剥き出しの文字を隣接パネルと並べると、下部の
    余白が「欠けている」ように見えるため、そうした箇所で使う。
    呼び出し側は戻り値のRegion内へadd_section_lead/add_item_list/
    add_paragraph_textbox等を配置する。
    """
    add_background_zone(slide, x, y, w, h, tone=tone, rounded=rounded)
    return Region(x, y, w, h).inset(inset_x, inset_y)


def add_card(slide, x, y, w, h, *, fill=PALETTE.surface_base,
             line=PALETTE.line_neutral, elevated=True):
    """独立した情報単位。既定で軽い影を付ける（elevated=Falseでフラットに）。"""
    return Box(slide, x, y, w, h, radius=LAYOUT.radius_base, fill=fill,
               line=line, line_width=Pt(0.7), elevated=elevated)


def add_focus_panel(slide, x, y, w, h, *, tone="solid"):
    if tone == "solid":
        fill, line = PALETTE.focus_primary, PALETTE.focus_primary
    elif tone == "brand":
        fill, line = PALETTE.surface_brand_soft, PALETTE.line_brand
    else:
        raise ValueError("toneは 'solid' または 'brand' を指定してください")
    return Box(slide, x, y, w, h, radius=LAYOUT.radius_base, fill=fill,
               line=line, line_width=Pt(1.0), elevated=False)


def add_key_message(slide, x, y, w, text, *, style="editorial"):
    """ページの結論・導入文を示す。

    styleは editorial/subtle/solid/card に加え、装飾なしの plain を選べる。
    plain は、ページタイトル直下の罫線に隣接して短い線を重ねたくない場合
    （例: ページ冒頭の導入文）に使う。
    """
    typography = _type_for(slide)
    heights = {"editorial": Inches(0.72), "subtle": Inches(0.7),
               "solid": Inches(0.78), "card": Inches(0.82),
               "plain": Inches(0.5)}
    h = heights.get(style, heights["editorial"])
    if style == "plain":
        return add_textbox(slide, x, y, w, h, text, size=typography.section,
                           bold=True, font=typography.headline_font)
    if style == "editorial":
        add_hairline(slide, x, y, Inches(0.38), color=PALETTE.blue, width=2)
        return add_textbox(slide, x, y + Inches(0.17), w, h - Inches(0.17), text,
                           size=typography.section, bold=True,
                           font=typography.headline_font)
    if style == "solid":
        box = Box(slide, x, y, w, h, radius=LAYOUT.radius_base,
                  fill=PALETTE.navy, line=None)
        color = PALETTE.white
    else:
        box = Box(slide, x, y, w, h, radius=LAYOUT.radius_base,
                  fill=PALETTE.grey_100, line=PALETTE.grey_300, line_width=Pt(0.7))
        color = PALETTE.ink
    tf = style_text_frame(box.text_frame, margin_x=Inches(0.18), margin_y=Inches(0.12),
                          vertical_anchor=MSO_ANCHOR.MIDDLE)
    run = tf.paragraphs[0].add_run(); run.text = text
    set_run(run, size=typography.body, color=color, bold=True,
            font=typography.body_font)
    return box


def add_item_list(slide, x, y, w, h, items, *, bullet="•", body_gap=3,
                  vertical_anchor=MSO_ANCHOR.TOP, adaptive=True):
    """複数項目を一つのtextboxとして配置し、手修正しやすく保つ。

    itemsは文字列、または {"title": ..., "body": ...} の配列。
    独立移動が必要な項目には使わない。section_leadの直下など見出しに
    連続させる場合はTOPのまま使う（既定）。項目数が領域に対して少ない
    場合、adaptive=True（既定）なら項目間の余白を上限付きで自動的に
    広げ、上詰めのまま下部の余白の割合を抑える。見出しを伴わない独立
    ブロック（カード、パネルなど）を領域全体で釣り合わせたい場合は
    vertical_anchor=MSO_ANCHOR.MIDDLEを指定する。

    実体はadd_text_list（marker="bullet"）。
    """
    return add_text_list(slide, x, y, w, h, items, marker="bullet",
                         bullet_char=bullet, gap=body_gap,
                         vertical_anchor=vertical_anchor, adaptive=adaptive)


def add_icon_list(slide, x, y, w, h, items, *, icon="check",
                  icon_color=PALETTE.blue, icon_size=None,
                  text_gap=Inches(0.16), body_gap=14):
    """アイコン付き箇条書きを描く。

    文章はadd_item_listと同じく一つのtextboxへまとめ、行ごとに別shapeへ
    分割しない（編集性を保つため）。アイコンは装飾図形であり文章そのもの
    ではないため、行ごとに独立した図形として置く。itemsは文字列の配列。
    iconは全行共通の名前、またはitemsと同数のnameリストを渡す
    （行ごとに変える場合）。

    icon_sizeを指定しない場合、本文フォントサイズの約2倍を目安に自動計算
    する。固定インチ値のままだとlarge-roomモードなど本文が大きいmodeで
    アイコンが相対的に小さく見える（business比で本文が約1.4倍でも
    アイコンは同じ大きさのまま、というズレが生じる）ため。

    実体はadd_text_list（marker="icon"）。
    """
    return add_text_list(slide, x, y, w, h, items, marker="icon", icon=icon,
                         icon_color=icon_color, icon_size=icon_size,
                         text_gap=text_gap, gap=body_gap)
