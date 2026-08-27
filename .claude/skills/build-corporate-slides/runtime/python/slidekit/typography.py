"""Atom層のうち、文字を扱う部分（図形はatoms.py、図形と文字にまたがる
部品はcomponents.py）。層構成の全体像はARCHITECTURE.mdを参照。

単独のテキスト描画（add_textbox / add_paragraph_textbox / set_run）と、
複数項目をまとめて1つのtextboxへ描くadd_text_list、テキストを含む定型
部品（Tag / Stat）が同居する。いずれもAtom層で、上下関係はない
（テキスト処理をこのファイルへ集約する方針による分け方）。
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .atoms import _filled_shape, add_hairline
from .icons import add_icon
from .textmetrics import (adaptive_gap_pt, char_width_factor,
                          estimate_item_list_height_pt,
                          estimate_paragraph_height_pt)
from .theme import PALETTE, TYPE


def _type_for(slide):
    slide_typography = getattr(slide, "_slidekit_typography", None)
    if slide_typography is not None:
        return slide_typography
    presentation = slide.part.package.presentation_part.presentation
    return getattr(presentation, "_slidekit_typography", TYPE)


def style_text_frame(text_frame, *, margin=0, margin_x=None, margin_y=None,
                     vertical_anchor=MSO_ANCHOR.TOP):
    """textboxの余白・折り返し・垂直配置を初期化する。

    margin_x/margin_yを指定すると左右/上下で異なる余白にできる（未指定
    の側はmargin、既定0を使う）。
    """
    text_frame.clear()
    text_frame.margin_left = text_frame.margin_right = (
        margin_x if margin_x is not None else margin)
    text_frame.margin_top = text_frame.margin_bottom = (
        margin_y if margin_y is not None else margin)
    text_frame.vertical_anchor = vertical_anchor
    text_frame.word_wrap = True
    return text_frame


def set_run(run, *, size=TYPE.body, color=PALETTE.ink, bold=False,
            font=TYPE.body_font):
    run.font.name = font
    properties = run._r.get_or_add_rPr()
    east_asian = properties.find(qn("a:ea"))
    if east_asian is None:
        east_asian = OxmlElement("a:ea")
        properties.append(east_asian)
    east_asian.set("typeface", font)
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return run


def _fill_text_frame(text_frame, text, *, size=TYPE.body, color=PALETTE.ink,
                     bold=False, font=TYPE.body_font, align=PP_ALIGN.LEFT,
                     margin=0, margin_x=None, margin_y=None, line_spacing=1.08,
                     vertical_anchor=MSO_ANCHOR.TOP):
    """既存のtext_frameへ単一run・単一段落のテキストを設定する。

    add_textbox（新規textboxを作ってから呼ぶ）とTag（既存図形のtext_frame
    へ直接呼ぶ）が共有する、テキスト設定だけの手順。
    """
    tf = style_text_frame(text_frame, margin=margin, margin_x=margin_x,
                          margin_y=margin_y, vertical_anchor=vertical_anchor)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.space_after = Pt(0)
    set_run(paragraph.add_run(), size=size, color=color, bold=bold, font=font)
    paragraph.runs[0].text = text
    return tf


def add_textbox(slide, x, y, w, h, text, *, size=TYPE.body,
                color=PALETTE.ink, bold=False, font=TYPE.body_font,
                align=PP_ALIGN.LEFT, margin=0, line_spacing=1.08):
    shape = slide.shapes.add_textbox(x, y, w, h)
    _fill_text_frame(shape.text_frame, text, size=size, color=color, bold=bold,
                     font=font, align=align, margin=margin, line_spacing=line_spacing)
    return shape


def add_paragraph_textbox(slide, x, y, w, h, paragraphs, *,
                          margin=0, vertical_anchor=MSO_ANCHOR.TOP):
    """複数の箇条書き・見出し・本文を一つの編集可能なtextboxへまとめる。

    paragraphsは {"segments": [(text, run_kwargs)], "space_after": 6,
    "level": 0, "bullet": False, "hanging_indent": Emu} の配列。意味上
    ひとまとまりの文章を、見た目上の行ごとに別shapeへ分割しないために使う。

    hanging_indentを指定すると、行頭の記号（箇条書きの記号や連番等、
    実体は先頭に付けた文字）の幅ぶんだけ、折り返した2行目以降を字下げ
    する。指定しないと記号の直後で折り返した際、2行目が左端から始まり
    記号に埋もれて見える。
    """
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = style_text_frame(shape.text_frame, margin=margin,
                          vertical_anchor=vertical_anchor)
    for index, spec in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = spec.get("align", PP_ALIGN.LEFT)
        paragraph.level = spec.get("level", 0)
        paragraph.line_spacing = spec.get("line_spacing", 1.08)
        paragraph.space_before = Pt(spec.get("space_before", 0))
        paragraph.space_after = Pt(spec.get("space_after", 0))
        if spec.get("bullet"):
            paragraph.text = "• "
        hanging_indent = spec.get("hanging_indent")
        if hanging_indent:
            pPr = paragraph._p.get_or_add_pPr()
            pPr.set("marL", str(int(hanging_indent)))
            pPr.set("indent", str(-int(hanging_indent)))
        for value, kwargs in spec.get("segments", []):
            run = paragraph.add_run()
            run.text = value
            set_run(run, **kwargs)
    return shape


def _normalize_list_item(item):
    if isinstance(item, str):
        return item, None
    return item.get("title", ""), item.get("body")


def _prefix_width(prefix, font_pt):
    """prefix文字列がfont_ptで描画される概算の幅（Emu）。行頭記号ぶんの
    折り返し2行目の字下げ幅として使う（実測レイアウトを持たないため、
    char_width_factorのヒューリスティックで概算する）。
    """
    return Pt(sum(char_width_factor(ch) for ch in prefix) * font_pt)


def _list_item_segments(marker, index, title, body, bullet_char, typography):
    """1項目分のsegments（(text, run_kwargs)の配列）と、行頭記号（連番・
    箇条書き記号）ぶんの折り返し字下げ幅（hanging_indent、Emu）を組み立てる。
    """
    segments = []
    hanging_indent = None
    if marker == "number":
        prefix = f"{index + 1:02}   "
        segments.append((prefix, {
            "size": typography.small, "color": PALETTE.blue,
            "bold": True, "font": typography.body_font,
        }))
        title_text, title_bold, indent = title, True, "      "
        hanging_indent = _prefix_width(prefix, typography.small.pt)
    elif marker == "bullet":
        prefix = f"{bullet_char}  "
        title_text, title_bold, indent = f"{prefix}{title}", True, "    "
        hanging_indent = _prefix_width(prefix, typography.body.pt)
    else:
        # "icon" / "none": 記号なしの本文（iconは呼び出し側が別図形で置く）
        title_text, title_bold, indent = title, False, "  "
    segments.append((title_text, {
        "size": typography.body, "color": PALETTE.text_primary,
        "bold": title_bold, "font": typography.body_font,
    }))
    if body:
        segments.append((f"\n{indent}{body}", {
            "size": typography.small, "color": PALETTE.text_secondary,
            "font": typography.body_font,
        }))
    return segments, hanging_indent


def add_text_list(slide, x, y, w, h, items, *, marker="bullet", bullet_char="・",
                  icon=None, icon_color=None, icon_size=None, text_gap=Inches(0.16),
                  divider=False, max_row_h=Inches(0.9),
                  vertical_anchor=MSO_ANCHOR.TOP, gap=3, adaptive=True):
    """複数項目のテキストを描く統合atom。

    旧add_item_list / add_icon_list / add_numbered_rowを1つの実装へ統合
    したもの。itemsは文字列、または{"title": ..., "body": ...}の配列。

    marker: "bullet"（既定、bullet_charでの箇条書き）/ "number"（01, 02...
      を自動採番）/ "icon"（iconで指定したアイコンを行ごとに独立した図形
      で置く）/ "none"（記号なし）。

    divider=False（既定）は全項目を1つのtextboxへまとめる。PowerPoint上
    での編集性を優先し、項目ごとにshapeへ分割しない。divider=Trueは各
    項目を独立した行として描き、行の間に実際の文字高さへ合わせた安全な
    罫線を引く（numbered_listのような単列の項目向け。安全な余白がない
    行は罫線を省略する）。項目数が少ないときはmax_row_hを上限に行の
    高さを詰め、領域内で上下中央へ配置する（上詰めのまま下部だけ余ると
    間延びして見えるため）。
    """
    typography = _type_for(slide)
    normalized = [_normalize_list_item(item) for item in items]

    if divider:
        row_h = Inches(0.9) if not items else min(max_row_h, h // len(items))
        block_h = row_h * len(items)
        cursor = y + max(0, (h - block_h) // 2)
        text_w_pt = w / 12700
        indent = "      " if marker == "number" else "    "
        for index, (title, body) in enumerate(normalized):
            segments, hanging_indent = _list_item_segments(
                marker, index, title, body, bullet_char, typography)
            add_paragraph_textbox(slide, x, cursor, w, Inches(0.76), [
                {"segments": segments, "space_after": 0,
                 "hanging_indent": hanging_indent}
            ])
            # 罫線は固定オフセットではなく実際の文字高さに合わせる。行が
            # 短い場合、固定オフセットのままだと罫線が自分の行から離れ、
            # 次の行の文字に近づきすぎて見える（区切り線が次行に属して
            # 見える）ため。row_hに対して安全な余白がない場合は、中途半端
            # な位置に引いて文字に被せるより罫線を省略する。
            content_pt = estimate_paragraph_height_pt(title, typography.body.pt, text_w_pt,
                                                       line_spacing=1.08)
            if body:
                content_pt += estimate_paragraph_height_pt(
                    indent + body, typography.small.pt, text_w_pt, line_spacing=1.08)
            content_h = Inches(content_pt / 72)
            offset = content_h + Inches(0.14)
            safe_offset = row_h - Inches(0.08)
            if safe_offset >= content_h + Inches(0.04):
                add_hairline(slide, x, cursor + min(offset, safe_offset), w)
            cursor += row_h
        return None

    if marker == "icon":
        names = icon if isinstance(icon, list) else [icon or "check"] * len(items)
        size = icon_size if icon_size is not None else Pt(typography.body.pt * 2.0)
        color = icon_color or PALETTE.blue
        text_x = x + size + text_gap
        text_w = w - size - text_gap
        text_w_pt = text_w / 12700
        heights_pt = [estimate_paragraph_height_pt(title, typography.body.pt, text_w_pt,
                                                    line_spacing=1.15)
                     for title, _ in normalized]
        gap_pt = adaptive_gap_pt(sum(heights_pt), len(items), h / 12700, base_gap=gap)
        paragraphs, cursor = [], y
        for index, ((title, body), name, item_h_pt) in enumerate(
                zip(normalized, names, heights_pt)):
            add_icon(slide, x, cursor, size, name, color=color)
            segments, _ = _list_item_segments(
                "icon", index, title, body, bullet_char, typography)
            paragraphs.append({
                "segments": segments,
                "line_spacing": 1.15,
                "space_after": gap_pt if index < len(items) - 1 else 0,
            })
            cursor += Inches(item_h_pt / 72) + Pt(gap_pt)
        return add_paragraph_textbox(slide, text_x, y, text_w, h, paragraphs)

    body_gap = gap
    if adaptive and items:
        content_pt = estimate_item_list_height_pt(typography, items, w / 12700, body_gap=0)
        body_gap = int(adaptive_gap_pt(content_pt, len(items), h / 12700, base_gap=gap))
    paragraphs = []
    for index, (title, body) in enumerate(normalized):
        segments, hanging_indent = _list_item_segments(
            marker, index, title, body, bullet_char, typography)
        paragraphs.append({
            "segments": segments,
            "space_after": body_gap if index < len(items) - 1 else 0,
            "hanging_indent": hanging_indent,
        })
    return add_paragraph_textbox(slide, x, y, w, h, paragraphs,
                                 vertical_anchor=vertical_anchor)


def Tag(slide, x, y, w, h, text, *, fill=None, line=PALETTE.grey_500,
        color=None, size=None, bold=True, pill=True):
    """ステータス・分類を示す小さなラベル（Atom層）。

    fill未指定（既定）は枠線のみのバッジ（表紙の開示区分など）。fillを
    指定すると単色塗りの状態チップになる（KPI/ステータス表示向け）。
    既定では角丸を高さの半分に固定し、完全な丸薬型にする。

    pill=Falseの場合は角のある矩形にする。表紙の開示区分バッジ
    （「部外秘」等）は社内文書の型として角ありで固定運用したいという
    実務要件のための例外で、それ以外の一般的なTagはpill=True（既定）
    のままにする。
    """
    typography = _type_for(slide)
    size = size or typography.small
    color = color or (PALETTE.white if fill else PALETTE.text_secondary)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if pill else MSO_SHAPE.RECTANGLE
    tag = _filled_shape(slide, shape_type, x, y, w, h,
                        radius=h / 2 if pill else None, fill=fill,
                        line=None if fill else line, line_width=Pt(1.0))
    _fill_text_frame(tag.text_frame, text, size=size, color=color, bold=bold,
                     font=typography.body_font, align=PP_ALIGN.CENTER,
                     margin_x=Inches(0.05), vertical_anchor=MSO_ANCHOR.MIDDLE)
    return tag


def Stat(slide, x, y, w, h, value, label=None, *, detail=None,
        value_size=None, value_color=PALETTE.navy,
        label_size=None, label_color=PALETTE.text_primary, label_bold=True,
        detail_color=PALETTE.text_secondary, vertical_anchor=MSO_ANCHOR.TOP):
    """数値＋ラベル（＋補足）を1つの編集可能テキストとして描く（Atom層）。

    stat_highlightの主指標・補足指標など、大きな数値を主役にする表現の
    最小単位。背景の面（Box/Card）は呼び出し側が別途用意する
    （Statはテキストだけを担当し、Box/Cardと自由に組み合わせられる
    ようにするため）。
    """
    typography = _type_for(slide)
    value_size = value_size or typography.metric
    label_size = label_size or typography.small
    paragraphs = [
        {"segments": [(value, {
            "size": value_size, "color": value_color, "bold": True,
            "font": typography.headline_font,
        })], "space_after": 4 if (label or detail) else 0},
    ]
    if label:
        paragraphs.append({"segments": [(label, {
            "size": label_size, "color": label_color, "bold": label_bold,
            "font": typography.body_font,
        })], "space_after": 4 if detail else 0})
    if detail:
        paragraphs.append({"segments": [(detail, {
            "size": typography.small, "color": detail_color,
            "font": typography.body_font,
        })]})
    return add_paragraph_textbox(slide, x, y, w, h, paragraphs,
                                 vertical_anchor=vertical_anchor)
