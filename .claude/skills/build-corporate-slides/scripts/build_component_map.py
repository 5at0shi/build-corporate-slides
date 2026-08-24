"""コンポーネント名称リファレンスを生成する（.claude/skills/build-corporate-slides/user-guide/component-map.pptx）。

実際のrenderer関数を直接呼び出して本物の出力を作り、その上に赤枠＋ラベルで
名称を注記する。座標は各rendererの内部計算と同じ`Region.rows/columns`を
使うため、注記位置は実際の描画と一致する。デザインシステムを変更したら
このスクリプトを再実行して参照を最新化する。

使い方（ワークスペース直下から）:
    ./.venv/bin/python .claude/skills/build-corporate-slides/scripts/build_component_map.py
    ./.venv/bin/python .claude/skills/build-corporate-slides/scripts/render_and_check.py \\
        .claude/skills/build-corporate-slides/user-guide/component-map.pptx \\
        --pdf .claude/skills/build-corporate-slides/user-guide/component-map.pdf
"""
import sys
from pathlib import Path

SKILL = Path(__file__).resolve().parents[1]
WORKSPACE_ROOT = Path.cwd()
sys.path.insert(0, str(SKILL / "runtime" / "python"))

from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from slidekit import DeckBuilder
from slidekit.components import (add_card, add_key_message,
                                 add_numbered_row, add_section_lead)
from slidekit.icons import add_icon
from slidekit.layout import Region, content_region
from slidekit.renderers import (render_comparison, render_evidence_and_decision,
                                render_table_with_conclusion)
from slidekit.theme import PALETTE

MARK = RGBColor(0xE8, 0x35, 0x35)


def annotate(slide, region, label, *, tag_align="left", pos="above"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, region.x, region.y,
                                 region.w, region.h)
    box.shadow.inherit = False
    box.fill.background()
    box.line.color.rgb = MARK
    box.line.width = Pt(1.25)

    tag_w = Inches(0.2 + 0.135 * len(label))
    if pos == "right":
        tag_x = region.x + region.w + Inches(0.1)
        tag_y = region.y + int(region.h / 2) - Inches(0.12)
    else:
        tag_x = region.x if tag_align == "left" else region.x + region.w - tag_w
        if pos == "inside":
            tag_y = region.y + Inches(0.04)
        else:
            tag_y = region.y - Inches(0.26)
            if tag_y < 0:
                tag_y = region.y + Inches(0.04)
    tag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tag_x, tag_y, tag_w,
                                 Inches(0.24))
    tag.shadow.inherit = False
    tag.fill.solid(); tag.fill.fore_color.rgb = MARK
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(9.5); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Hiragino Sans"


def note(slide, x, y, w, text):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.5))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = Pt(10.5); run.font.italic = True
    run.font.color.rgb = PALETTE.text_secondary
    run.font.name = "Hiragino Sans"


builder = DeckBuilder.from_workspace(WORKSPACE_ROOT)

# ---------------------------------------------------------------- 1. Cover
builder.add_cover(
    "生成AIの社内導入に向けたPoC計画",
    subtitle="安全に、効果を測りながら進める10週間の実証",
    eyebrow="AI ADOPTION STRATEGY",
    brand_side="right", brand_shape="diagonal")
slide = builder.presentation.slides[-1]
content_x, content_w = Inches(0.67), Inches(8.5)
annotate(slide, Region(content_x, Inches(2.08), content_w, Inches(0.28)),
         "Eyebrow")
annotate(slide, Region(content_x, Inches(2.48), content_w, Inches(1.3)),
         "Title", pos="inside")
annotate(slide, Region(content_x, Inches(4.03), content_w, Inches(0.62)),
         "Subtitle")
annotate(slide, Region(Inches(9.93), Inches(0), Inches(3.4), Inches(7.5)),
         "Cover Brand Field", pos="inside")
annotate(slide, Region(Inches(11.42), Inches(0.32), Inches(1.22), Inches(0.3)),
         "Badge", pos="inside")
annotate(slide, Region(Inches(9.2), Inches(0.72), Inches(3.45), Inches(0.52)),
         "Meta Info")
# ロゴの実際の位置はadd_cover内部で画像の高さから動的に計算されるため、
# 座標を推測せず、実際に配置されたPicture図形から直接取得する。
logo_shape = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
annotate(slide, Region(logo_shape.left, logo_shape.top, logo_shape.width,
                       logo_shape.height), "Logo", tag_align="right")

# ---------------------------------------------------- 2. comparison(asymmetric)
spec2 = {
    "title": "外部環境の変化と社内の課題が重なっている", "density": "standard",
    "variant": "asymmetric",
    "primary_message": "PoCを通じて、会社としての生成AI活用方針を早期に確立する必要がある",
    "left": {"heading": "外部環境の変化", "items": [
        {"title": "競合他社での活用進展", "body": "業務プロセスへの生成AI組み込みが加速している"},
        {"title": "人材獲得競争の激化", "body": "先進的な業務環境が採用力に影響し始めている"}]},
    "right": {"heading": "社内の課題", "items": [
        {"title": "定型業務の負荷", "body": "資料作成や情報整理に多くの時間が割かれている"},
        {"title": "利用の属人化", "body": "一部社員が個人判断でツールを使い始めている"}]},
}
render_comparison(builder, spec2, 2)
slide = builder.presentation.slides[-1]
area = content_region()
body, conclusion = area.rows([4.35, 0.72], gap=Inches(0.24))
left, right = body.columns([1.08, 0.92], gap="wide")
annotate(slide, Region(Inches(0.67), Inches(0.62), Inches(11.99), Inches(0.52)), "Title")
annotate(slide, Region(left.x, left.y, left.w, Inches(0.42)), "Section Lead")
annotate(slide, Region(left.x, left.y + Inches(0.62), left.w, left.h - Inches(0.68)),
         "Item List")
annotate(slide, Region(right.x - Inches(0.18), right.y - Inches(0.08),
                           right.w + Inches(0.36), right.h + Inches(0.12)),
         "Background Zone (neutral)", tag_align="right")
annotate(slide, Region(conclusion.x, conclusion.y, conclusion.w, conclusion.h),
         "Key Message (subtle)")
note(slide, Inches(0.67), Inches(7.05), Inches(11), "赤枠は名称の注記。本編デザインには含まれない。")

# ---------------------------------------------- 3. evidence_and_decision
spec3 = {
    "title": "全社一斉導入ではなく、小さく検証してから広げる", "density": "standard",
    "primary_message": "PoCで安全性と効果を確認し、本格導入の判断材料をそろえる",
    "evidence_heading": "全社一斉導入を避ける理由",
    "evidence": [
        {"title": "リスクの見極めが不十分", "body": "情報管理や誤利用への対応が定まっていない"},
        {"title": "効果が未検証", "body": "どの業務でどれだけの効果が出るか分からない"}],
    "decision_heading": "推奨する進め方",
    "decision_detail": "対象業務を絞ったPoCを10週間実施し、判断材料をそろえる。",
}
render_evidence_and_decision(builder, spec3, 3)
slide = builder.presentation.slides[-1]
left3, right3 = area.columns([1.55, 1], gap="wide")
annotate(slide, Region(left3.x - Inches(0.12), left3.y - Inches(0.06),
                           left3.w + Inches(0.24), left3.h + Inches(0.12)),
         "Panel (background zone: neutral)")
annotate(slide, Region(right3.x, right3.y, right3.w, right3.h),
         "Panel (background zone: brand-soft)", tag_align="right")
note(slide, right3.x, right3.y + right3.h + Inches(0.08), right3.w,
     "中の見出し・結論文はKey Messageと同じ役割")

# --------------------------------------------------- 4. table_with_conclusion
spec4 = {
    "title": "定量・定性の両面から効果を測定する", "density": "dense",
    "primary_message": "必須項目をすべて満たした場合に、本格導入の検討へ進める",
    "columns": [
        {"key": "criterion", "label": "評価項目", "weight": 1.6},
        {"key": "measure", "label": "確認方法", "weight": 2.6},
        {"key": "decision", "label": "判定", "weight": 1, "align": "center"}],
    "rows": [
        {"criterion": "作業時間削減", "measure": "現行業務との所要時間比較", "decision": "必須",
         "_highlight": True},
        {"criterion": "品質維持", "measure": "担当者レビューと修正量", "decision": "必須"},
        {"criterion": "安全性", "measure": "情報漏えい・誤送信の有無", "decision": "必須"}],
}
render_table_with_conclusion(builder, spec4, 4)
slide = builder.presentation.slides[-1]
table_region, conclusion4 = area.rows([4.35, 0.72], gap=Inches(0.24))
row_h = int(table_region.h / 4)
annotate(slide, Region(table_region.x, table_region.y, table_region.w, row_h),
         "Table Header Row")
annotate(slide, Region(table_region.x, table_region.y + row_h, table_region.w, row_h),
         "Highlighted Row（例外行）")
annotate(slide, Region(conclusion4.x, conclusion4.y, conclusion4.w, conclusion4.h),
         "Key Message (subtle)")

# ------------------------------------------------------- 5. composite parts
slide = builder.add_slide("その他のパーツ名称", density="standard", page=5)[0]
area5 = content_region()
card_row = area5.rows([2.2, 1.3, 0.9, 0.62], gap=Inches(0.2))
cards = card_row[0].columns([1, 1, 1], gap="standard")
for index, region in enumerate(cards):
    add_card(slide, region.x, region.y, region.w, region.h)
    add_icon(slide, region.x + Inches(0.25), region.y + Inches(0.25), Inches(0.5),
             ["check", "warning", "target"][index], color=PALETTE.blue)
    annotate(slide, region, "Card" if index == 0 else "", tag_align="left")
annotate(slide, Region(cards[0].x + Inches(0.25), cards[0].y + Inches(0.25),
                            Inches(0.5), Inches(0.5)), "Icon", tag_align="right")

row_region = card_row[1]
add_numbered_row(slide, row_region.x, row_region.y, row_region.w, 1,
                 "予算承認", "ツール利用料・導入支援に関わる費用")
annotate(slide, Region(row_region.x, row_region.y, row_region.w, Inches(0.9)),
         "Numbered Row")

badge_region = card_row[2]
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_region.x, badge_region.y,
                               Inches(1.3), Inches(0.32))
badge.shadow.inherit = False
badge.fill.background(); badge.line.color.rgb = PALETTE.grey_500; badge.line.width = Pt(1.0)
tf = badge.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "サンプル"; r.font.size = Pt(10); r.font.bold = True
r.font.color.rgb = PALETTE.text_secondary; r.font.name = "Hiragino Sans"
annotate(slide, Region(badge_region.x, badge_region.y, Inches(1.3), Inches(0.32)),
         "Badge", pos="right")

add_key_message(slide, card_row[3].x, card_row[3].y, card_row[3].w,
                "強い結論を示すときはsolidを使う", style="solid")
annotate(slide, card_row[3], "Key Message (solid)", tag_align="right")

output = builder.save(SKILL / "user-guide" / "component-map.pptx")
print("PPTX:", output)
