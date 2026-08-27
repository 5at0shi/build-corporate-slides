"""外部アセットに依存しない、ベクター描画のみのアイコンセット。

ライセンス上の懸念を避けるため、既存のアイコンフォント・SVGセットは使わず、
python-pptxの図形（Oval/RoundedRectangle/Freeform）だけで構成する。
座標はすべて0.0〜1.0のローカル単位（アイコンの正方形キャンバスに対する
比率）で書き、`add_icon`が実寸へスケール・配置する。

線（チェックマーク・時計の針）は、開いたfreeformのストロークではなく、
細い RoundedRectangle を回転させて作る（開いたfreeformは一部レンダラーで
意図通り描画されないため）。回転は図形自身の中心を軸に回るため、始点を
固定するための補正（_hand内のax/ay計算）を行っている。
"""
import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .atoms import _flat
from .theme import PALETTE

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ICON_NAMES = {"check", "warning", "people", "document", "clock", "target",
             "growth", "cost", "calendar", "idea", "location", "building",
             "globe", "shield"}


def _oval(slide, x0, y0, x1, y1, *, x, y, size, fill=None, line=None,
          line_width=Pt(1.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + int(x0 * size), y + int(y0 * size),
        int((x1 - x0) * size), int((y1 - y0) * size))
    _flat(shape)
    if fill is not None:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line; shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _rounded_rect(slide, x0, y0, x1, y1, *, x, y, size, fill, rounding=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + int(x0 * size), y + int(y0 * size),
        int((x1 - x0) * size), int((y1 - y0) * size))
    _flat(shape, rounding=rounding)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _preset(slide, shape_type, x0, y0, x1, y1, *, x, y, size, fill, rounding=None):
    """OVAL/ROUNDED_RECTANGLE以外の定型オートシェイプ用（三角形・台形等）。

    自由形状（freeform）はこの環境のレンダリングパイプラインで意図通りに
    描画されないことがあるため、アイコンでは定型プリセットだけを使う。
    roundingはROUND_2_SAME_RECTANGLE等、角丸系のadjustmentを持つ形状用。
    """
    shape = slide.shapes.add_shape(
        shape_type, x + int(x0 * size), y + int(y0 * size),
        int((x1 - x0) * size), int((y1 - y0) * size))
    _flat(shape, rounding=rounding)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _hand_between(slide, p0, p1, thickness, *, x, y, size, color):
    """p0からp1へ向かう、太さthicknessの棒を描く（p0を固定端として）。

    RoundedRectangleを回転させて作る。回転は図形自身の中心を軸に回るため、
    p0が回転後も動かないよう、開始位置をあらかじめ補正している。
    """
    dx, dy = p1[0] - p0[0], p1[1] - p0[1]
    length = math.hypot(dx, dy)
    ux, uy = dx / length, dy / length
    angle_deg = math.degrees(math.atan2(ux, -uy))
    theta = math.radians(angle_deg)
    anchor_x = p0[0] + (length / 2) * math.sin(theta)
    anchor_y = p0[1] + (length / 2) * (1 - math.cos(theta))

    w = max(1, int(thickness * size))
    h = max(1, int(length * size))
    left = x + int(anchor_x * size) - w // 2
    top = y + int(anchor_y * size) - h
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    _flat(shape)
    if len(shape.adjustments):
        shape.adjustments[0] = 0.5
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = angle_deg
    return shape


def _draw_check(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _hand_between(slide, (0.26, 0.5), (0.43, 0.68), 0.1,
                 x=x, y=y, size=size, color=_WHITE)
    _hand_between(slide, (0.43, 0.68), (0.77, 0.3), 0.1,
                 x=x, y=y, size=size, color=_WHITE)
    # 2本のhand_betweenは同じ頂点(0.43, 0.68)を共有するが、回転計算の
    # 誤差でわずかにずれて継ぎ目が荒く見えるため、同じ半径の円を重ねて
    # 継ぎ目を確実に滑らかな丸みへ揃える。
    _oval(slide, 0.38, 0.63, 0.48, 0.73, x=x, y=y, size=size, fill=_WHITE)


def _draw_warning(slide, x, y, size, color):
    # 他の5種と同じ「色チップ＋白い図形」の言語に揃える
    # （三角形そのものを外形にすると、丸チップの並びの中で浮いて見えるため）。
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _preset(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, 0.2, 0.2, 0.8, 0.74,
           x=x, y=y, size=size, fill=_WHITE)
    _rounded_rect(slide, 0.465, 0.36, 0.535, 0.56, x=x, y=y, size=size,
                  fill=color, rounding=0.5)
    _oval(slide, 0.465, 0.62, 0.535, 0.67, x=x, y=y, size=size, fill=color)


def _draw_people(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _oval(slide, 0.36, 0.16, 0.64, 0.44, x=x, y=y, size=size, fill=_WHITE)
    _preset(slide, MSO_SHAPE.TRAPEZOID, 0.26, 0.58, 0.74, 0.92,
           x=x, y=y, size=size, fill=_WHITE)


def _draw_document(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _rounded_rect(slide, 0.3, 0.18, 0.7, 0.82, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.12)
    for line_y in (0.36, 0.48, 0.6):
        _rounded_rect(slide, 0.4, line_y, 0.6, line_y + 0.045, x=x, y=y,
                      size=size, fill=color, rounding=0.5)


def _draw_clock(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _oval(slide, 0.16, 0.16, 0.84, 0.84, x=x, y=y, size=size,
          line=_WHITE, line_width=Pt(size.pt * 0.07))
    _hand_between(slide, (0.5, 0.5), (0.5, 0.26), 0.08,
                 x=x, y=y, size=size, color=_WHITE)
    _hand_between(slide, (0.5, 0.5), (0.68, 0.42), 0.08,
                 x=x, y=y, size=size, color=_WHITE)
    _oval(slide, 0.46, 0.46, 0.54, 0.54, x=x, y=y, size=size, fill=_WHITE)


def _draw_target(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _oval(slide, 0.18, 0.18, 0.82, 0.82, x=x, y=y, size=size, fill=_WHITE)
    _oval(slide, 0.36, 0.36, 0.64, 0.64, x=x, y=y, size=size, fill=color)


def _draw_growth(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    for bar_x, top in ((0.27, 0.58), (0.44, 0.46), (0.61, 0.32)):
        _rounded_rect(slide, bar_x, top, bar_x + 0.12, 0.72, x=x, y=y,
                      size=size, fill=_WHITE, rounding=0.25)


def _draw_cost(slide, x, y, size, color):
    # 円+横棒1本は「禁止」記号と紛らわしいため、コインを横から見た
    # 積み重ねで表す（白い楕円を3枚、色の輪郭線で個々の縁を見せる）。
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    for coin_y in (0.56, 0.46, 0.36):
        _oval(slide, 0.24, coin_y, 0.76, coin_y + 0.16, x=x, y=y, size=size,
              fill=_WHITE, line=color, line_width=Pt(size.pt * 0.045))


def _draw_calendar(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    # リング（背景円の上に重なる）を先に描き、本体で下端を隠して繋がって
    # 見えるようにする。
    _rounded_rect(slide, 0.32, 0.14, 0.4, 0.3, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.5)
    _rounded_rect(slide, 0.6, 0.14, 0.68, 0.3, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.5)
    # 本体は常に白一色にする。ヘッダーを背景円と同色で塗ると境界が同化
    # して輪郭が読めなくなるため（実際に起きていた不具合）、ヘッダーは
    # 本体内部の細い罫線として表現し、外形は必ず白のまま保つ。
    _rounded_rect(slide, 0.2, 0.24, 0.8, 0.82, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.12)
    _rounded_rect(slide, 0.26, 0.36, 0.74, 0.4, x=x, y=y, size=size,
                  fill=color, rounding=0.5)
    # 日付の格子（2行×3列）を足し、単なるカードと区別できるようにする。
    for row_y in (0.5, 0.64):
        for col_x in (0.28, 0.44, 0.6):
            _rounded_rect(slide, col_x, row_y, col_x + 0.12, row_y + 0.09,
                         x=x, y=y, size=size, fill=color, rounding=0.3)


def _draw_idea(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _oval(slide, 0.28, 0.14, 0.72, 0.58, x=x, y=y, size=size, fill=_WHITE)
    _rounded_rect(slide, 0.4, 0.54, 0.6, 0.72, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.3)
    _rounded_rect(slide, 0.43, 0.74, 0.57, 0.8, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.5)


def _draw_location(slide, x, y, size, color):
    # 円と三角形の重なりが浅いと分離した「風船」に見えるため、同じ幅で
    # 深く重ね、先端まで伸ばして1本のピン形状に見せる。
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _oval(slide, 0.24, 0.14, 0.76, 0.66, x=x, y=y, size=size, fill=_WHITE)
    _preset(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, 0.24, 0.42, 0.76, 0.88,
           x=x, y=y, size=size, fill=_WHITE)
    _oval(slide, 0.42, 0.28, 0.58, 0.44, x=x, y=y, size=size, fill=color)


def _draw_building(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _rounded_rect(slide, 0.26, 0.18, 0.74, 0.82, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.08)
    for row_y in (0.3, 0.48, 0.66):
        for col_x in (0.35, 0.565):
            _rounded_rect(slide, col_x, row_y, col_x + 0.08, row_y + 0.1,
                          x=x, y=y, size=size, fill=color, rounding=0.2)


def _draw_globe(slide, x, y, size, color):
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    _oval(slide, 0.16, 0.16, 0.84, 0.84, x=x, y=y, size=size,
          line=_WHITE, line_width=Pt(size.pt * 0.07))
    _rounded_rect(slide, 0.16, 0.465, 0.84, 0.535, x=x, y=y, size=size,
                  fill=_WHITE, rounding=0.5)
    _oval(slide, 0.35, 0.16, 0.65, 0.84, x=x, y=y, size=size,
          line=_WHITE, line_width=Pt(size.pt * 0.06))


def _draw_shield(slide, x, y, size, color):
    # 正五角形はキャンバス中央(0.14-0.82)に置くと、下端の尖りぶんだけ
    # 視覚的な重心が上寄りに感じられる（先端は面積が小さく「軽い」ため）。
    # 上下0.04ぶん下へずらし、キャンバス内でのバランスを取る。
    _oval(slide, 0, 0, 1, 1, x=x, y=y, size=size, fill=color)
    shield = _preset(slide, MSO_SHAPE.REGULAR_PENTAGON, 0.24, 0.18, 0.76, 0.86,
                     x=x, y=y, size=size, fill=_WHITE)
    shield.rotation = 180
    _hand_between(slide, (0.4, 0.52), (0.48, 0.62), 0.09,
                 x=x, y=y, size=size, color=color)
    _hand_between(slide, (0.48, 0.62), (0.64, 0.4), 0.09,
                 x=x, y=y, size=size, color=color)
    # 2本のhand_betweenの継ぎ目を、同半径の円で確実に滑らかにする
    # （checkと同じ理由）。
    _oval(slide, 0.435, 0.575, 0.525, 0.665, x=x, y=y, size=size, fill=color)


_DRAWERS = {
    "check": _draw_check,
    "warning": _draw_warning,
    "people": _draw_people,
    "document": _draw_document,
    "clock": _draw_clock,
    "target": _draw_target,
    "growth": _draw_growth,
    "cost": _draw_cost,
    "calendar": _draw_calendar,
    "idea": _draw_idea,
    "location": _draw_location,
    "building": _draw_building,
    "globe": _draw_globe,
    "shield": _draw_shield,
}


def add_icon(slide, x, y, size, name, *, color=PALETTE.blue):
    """アイコンを(x, y)を左上、一辺sizeの正方形として描く。

    nameは check/warning/people/document/clock/target/growth/cost/calendar/
    idea/location/building/globe/shield のいずれか。
    いずれも「colorの円チップ＋白い図形」で構成し、どんな背景の上でも
    視認できるようにしている。装飾ではなく、状態・分類を示す場合に使う。
    """
    if name not in _DRAWERS:
        raise ValueError(f"未定義のiconです: {name} (候補: {sorted(ICON_NAMES)})")
    _DRAWERS[name](slide, x, y, size, color)
