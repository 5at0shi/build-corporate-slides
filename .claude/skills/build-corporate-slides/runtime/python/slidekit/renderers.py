"""Renderer層: Layout/Atom/Fragmentを組み合わせ、意味を持った1ページを
完成させる最上層。

ここで初めてビジネス語彙（「比較」「結論」「ゲート」等）が現れる。1関数=
1renderer type=YAMLのtype:フィールド1つに対応し、これがユーザー/AIが
実際にcontent-authoringで触るインターフェース（renderer-catalog.md）。
表紙・章扉・ページヘッダーも同じRenderer層だが、ページ外枠を描くため
pageframe.pyに分け、DeckBuilder経由で呼ぶ（循環importを避けるため）。
層構成の全体像はARCHITECTURE.mdを参照。
"""
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .atoms import Box, Connector, add_hairline
from .builder import DeckBuilder
from .charts import add_native_chart
from .components import (SECTION_LEAD_GAP, add_background_zone,
                         add_emphasis_zone, add_item_list, add_key_message,
                         add_panel, add_section_lead)
from .fragments import BoxGrid, MarkerOverlay, ProportionalStack, RadialLayout
from .layout import Region
from .preflight import require_valid_content
from .images import add_image_contain
from .tables import add_data_table
from .textmetrics import (adaptive_gap_pt, centered_gap_pt,
                          char_width_factor,
                          estimate_item_list_height_pt, estimate_line_count,
                          estimate_paragraph_height_pt)
from .theme import PALETTE
from .typography import (Stat, _type_for, add_paragraph_textbox,
                         add_text_list, add_textbox)


HEADING_BLOCK_H = Inches(0.62)


def _items(section):
    if isinstance(section, list):
        return section
    return section.get("items", [])


def _heading(section, fallback):
    return section.get("heading", fallback) if isinstance(section, dict) else fallback


def _adaptive_gap(items, available_h, *, base_gap, per_item_pt):
    if not items:
        return base_gap
    content_pt = per_item_pt * len(items)
    return int(adaptive_gap_pt(content_pt, len(items), available_h / 12700,
                               base_gap=base_gap))


def _lead_list(slide, region, heading, items, *, color=PALETTE.line_brand,
              bullet="•", bottom_pad=Inches(0.06)):
    """見出し(add_section_lead)＋その下のリスト(add_item_list)という、
    複数rendererで繰り返される組み合わせをまとめる。リストはHEADING_BLOCK_H
    ぶん見出しの下から始まる。
    """
    add_section_lead(slide, region.x, region.y, region.w, heading, color=color)
    add_item_list(slide, region.x, region.y + HEADING_BLOCK_H, region.w,
                  region.h - HEADING_BLOCK_H - bottom_pad, items, bullet=bullet)


_TONE_COLORS = {"positive": PALETTE.positive, "negative": PALETTE.negative,
                "warning": PALETTE.warning}


def _tone_color(item, default):
    """item["tone"]（positive/negative/warning）を符号専用の色へ変換する。

    数値の先頭が"-"かどうかでは判定しない（削減率のように、マイナスの
    数値が良い結果を意味することがあるため）。toneは呼び出し側（content）
    が意味を判断して明示する。
    """
    return _TONE_COLORS.get(item.get("tone"), default)


# YAMLのtone名（stat_highlight等と同じ語彙）から、add_background_zoneの
# tone名への対応。未定義の名前は呼び出し側が既定値へ落とす（生成を
# 止めない。誤った色より、生成が止まってページごと失われる方が痛い）。
_ZONE_TONES = {
    "brand": "brand-soft", "teal": "teal-soft", "neutral": "neutral",
    "positive": "positive-soft", "negative": "negative-soft",
    "warning": "warning-soft",
}


def _number(value):
    """YAMLの値を数値へ寄せる。数値化できない値は0として扱う。

    数量を描くrenderer（waterfall）は、値が文字列や空でも生成そのものを
    止めない。0として描けば「値が入っていない」ことがページ上で見える。
    """
    try:
        return float(value)
    except (TypeError, ValueError):
        return 0.0


def _delta_label(value, kind):
    """waterfallの値ラベルの既定表記。増減には符号を明示する。"""
    number = int(value) if float(value).is_integer() else value
    return f"{number:,}" if kind == "total" else f"{number:+,}"


def _clamp_period(value, count):
    """timelineのstart/endを1〜countへ丸める。

    periodsの範囲外を書いてもページを落とさず、軸の端に寄せて描く
    （範囲外であること自体は目で見て分かる）。
    """
    try:
        index = int(value)
    except (TypeError, ValueError):
        index = 1
    return max(1, min(count, index))


def _conclude(slide, conclusion, spec, *, default_style="subtle"):
    """rendererの末尾で結論(primary_message)を示す定型パターン。
    add_key_messageへ橋渡しするだけで、renderer固有の意味は持たない。

    見た目の指定はmessage_style（primary_messageの見た目を決めるため、
    numbered_listのtop位置=導入文でも同じ名前で通るこの名前に統一する）。
    """
    return add_key_message(slide, conclusion.x, conclusion.y, conclusion.w,
                           spec["primary_message"],
                           style=spec.get("message_style", default_style))


def render_cover(builder, spec, page):
    brand_width = spec.get("brand_width")
    builder.add_cover(
        spec["title"],
        subtitle=spec.get("subtitle"),
        eyebrow=spec.get("eyebrow"),
        brand_side=spec.get("brand_side", "right"),
        brand_shape=spec.get("brand_shape", "diagonal"),
        brand_width=Inches(brand_width) if brand_width else None,
        classification=spec.get("classification"),
        created=spec.get("created"),
    )


def render_comparison(builder, spec, page):
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    body, conclusion = area.rows([4.35, 0.72], gap=Inches(0.24))
    variant = spec.get("variant", "balanced")
    weights = [1, 1] if variant == "balanced" else [1.08, 0.92]
    left, right = body.columns(weights, gap="wide")
    if variant == "asymmetric":
        add_emphasis_zone(slide, right, tone="neutral")
    left_items = _items(spec.get("left", {}))
    right_items = _items(spec.get("right", {}))
    _lead_list(slide, left, _heading(spec.get("left", {}), "左側"), left_items)
    _lead_list(slide, right, _heading(spec.get("right", {}), "右側"), right_items,
              color=PALETTE.accent_secondary, bullet="—")
    _conclude(slide, conclusion, spec)


def render_evidence_and_decision(builder, spec, page):
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    left, right = area.columns([1.55, 1], gap="wide")
    evidence = spec.get("evidence", [])

    # 右の推奨パネルと同じ「意図的な余白を持つ面」に見せるため、左も淡色の
    # 面で包み、見出し+リストのブロックをその中で釣り合わせる。背景のない
    # 剥き出しの文字のまま下部が空くと、右のパネルとの対比で「欠けている」
    # ように見えるため。
    left_inner = add_panel(slide, left.x - Inches(0.12), left.y - Inches(0.06),
                           left.w + Inches(0.24), left.h + Inches(0.12),
                           tone="neutral")
    # 項目間隔を先に決めてからブロックの高さを出し、その高さで中央寄せする。
    # 最小間隔のまま中央へ置くと、周囲の余白だけが広がって項目群が窮屈に
    # 固まって見えるため（visual-quality.md）。
    content_pt = estimate_item_list_height_pt(
        typography, evidence, left_inner.w / 12700, body_gap=0,
        title_prefix="•  ")
    gap_pt = centered_gap_pt(content_pt, len(evidence),
                             (left_inner.h - HEADING_BLOCK_H) / 12700)
    list_pt = content_pt + gap_pt * max(0, len(evidence) - 1)
    list_h = min(left_inner.h - HEADING_BLOCK_H, Inches(list_pt / 72))
    top = left_inner.y + max(0, (left_inner.h - HEADING_BLOCK_H - list_h) // 2)
    add_section_lead(slide, left_inner.x, top, left_inner.w,
                     spec.get("evidence_heading", "判断の根拠"))
    add_item_list(slide, left_inner.x, top + HEADING_BLOCK_H, left_inner.w,
                  list_h, evidence, body_gap=int(gap_pt), adaptive=False)

    inner = add_panel(slide, right.x, right.y, right.w, right.h,
                      tone="brand-soft", inset_x=Inches(0.34), inset_y=Inches(0.3))
    paragraphs = [
        {"segments": [(spec.get("decision_heading", "推奨方針"), {
            "size": typography.small, "color": PALETTE.blue, "bold": True,
            "font": typography.body_font,
        })], "space_after": 10},
        {"segments": [(spec["primary_message"], {
            "size": typography.section, "color": PALETTE.text_primary,
            "bold": True, "font": typography.headline_font,
        })], "space_after": 14 if spec.get("decision_detail") else 0,
         "line_spacing": 1.05},
    ]
    if spec.get("decision_detail"):
        paragraphs.append({"segments": [(spec["decision_detail"], {
            "size": typography.body, "color": PALETTE.text_secondary,
            "font": typography.body_font,
        })], "line_spacing": 1.15})
    add_paragraph_textbox(slide, inner.x, inner.y, inner.w, inner.h,
                          paragraphs, vertical_anchor=MSO_ANCHOR.MIDDLE)


def render_scope_and_exclusions(builder, spec, page):
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    scope, exclusions = area.columns([2.2, 1], gap="wide")
    add_background_zone(slide, scope.x, scope.y, scope.w, scope.h,
                        tone="brand-soft", rounded=True)
    scope_inner = scope.inset(Inches(0.32), Inches(0.32))
    items = spec.get("scope", [])
    columns = scope_inner.columns([1] * max(1, len(items)), gap="standard")

    # 見出し＋項目を1つのブロックとして、period行の上までの帯の中で上下中央
    # へ置く。固定オフセットのままだとゾーン下部だけが大きく空いて間延びして
    # 見える（visual-quality.md）。見出しだけ上に残して項目を動かすと見出しと
    # 項目の結びつきが切れるため、evidence_and_decisionと同じくブロックごと
    # 動かす。
    typography = _type_for(slide)
    heading_block_h = Inches(0.78)
    label_block_h = Inches(0.44)
    period_reserve = Inches(0.62) if spec.get("period") else Inches(0)
    band_h = max(0, scope_inner.h - period_reserve)
    column_w_pt = (columns[0].w / 12700) if columns else 0
    content_pt = 0.0
    for item in items:
        height = estimate_paragraph_height_pt(
            item.get("title", ""), typography.section.pt, column_w_pt,
            line_spacing=1.08, space_after=7)
        height += estimate_paragraph_height_pt(
            item.get("body", ""), typography.small.pt, column_w_pt,
            line_spacing=1.08)
        content_pt = max(content_pt, height)
    block_h = heading_block_h + label_block_h + Inches(content_pt / 72)
    block_top = scope_inner.y + max(0, (band_h - block_h) // 2)
    label_y = block_top + heading_block_h
    text_y = label_y + label_block_h

    add_section_lead(slide, scope_inner.x, block_top, scope_inner.w,
                     spec.get("scope_heading", "対象範囲"))
    for index, item in enumerate(items):
        region = columns[index]
        label = item.get("label", f"{index + 1:02}")
        add_textbox(slide, region.x, label_y, region.w,
                    Inches(0.28), label, size=_type_for(slide).small,
                    color=PALETTE.blue, bold=True)
        add_paragraph_textbox(slide, region.x, text_y,
                              region.w, max(Inches(0.4), band_h - label_block_h), [
            {"segments": [(item.get("title", ""), {
                "size": _type_for(slide).section,
                "color": PALETTE.text_primary,
                "bold": True,
                "font": _type_for(slide).body_font,
            })], "space_after": 7},
            {"segments": [(item.get("body", ""), {
                "size": _type_for(slide).small,
                "color": PALETTE.text_secondary,
                "font": _type_for(slide).body_font,
            })]},
        ])
    if spec.get("period"):
        add_textbox(slide, scope_inner.x, scope_inner.y + scope_inner.h - Inches(0.42),
                    scope_inner.w, Inches(0.28), spec["period"],
                    size=_type_for(slide).small, color=PALETTE.text_primary,
                    bold=True)

    add_section_lead(slide, exclusions.x, exclusions.y, exclusions.w,
                     spec.get("exclusions_heading", "対象外"),
                     color=PALETTE.text_secondary)
    exclusions_items = spec.get("exclusions", [])
    add_item_list(slide, exclusions.x, exclusions.y + Inches(0.7),
                  exclusions.w, exclusions.h - Inches(1.35),
                  exclusions_items, bullet="—", body_gap=7)
    add_textbox(slide, exclusions.x, exclusions.y + exclusions.h - Inches(0.42),
                exclusions.w, Inches(0.34), spec["primary_message"],
                size=_type_for(slide).small, color=PALETTE.text_primary, bold=True)


def render_process_with_gates(builder, spec, page):
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    phase_row, work_row, gate_row, conclusion = area.rows(
        [0.72, 2.35, 1.1, 0.62], gap=Inches(0.19))
    phases = spec.get("phases", [])
    phase_weights = [phase.get("weight", 1) for phase in phases] or [1]
    columns = phase_row.columns(phase_weights, gap="tight")
    tones = ["brand-soft", "neutral", "teal-soft"]
    for index, (phase, region) in enumerate(zip(phases, columns)):
        add_background_zone(slide, region.x, region.y, region.w, region.h,
                            tone=tones[index % len(tones)], rounded=True)
        add_textbox(slide, region.x + Inches(0.15), region.y + Inches(0.15),
                    region.w - Inches(0.3), Inches(0.3), phase.get("title", ""),
                    size=_type_for(slide).body, color=PALETTE.text_primary,
                    bold=True, align=PP_ALIGN.CENTER)

    work_columns = work_row.columns(phase_weights, gap="tight")
    for index, (phase, region) in enumerate(zip(phases, work_columns)):
        add_textbox(slide, region.x, region.y + Inches(0.18), region.w,
                    Inches(0.28), phase.get("label", f"{index + 1:02}"),
                    size=_type_for(slide).small, color=PALETTE.blue, bold=True)
        # phasesは最大6分割まで想定する狭い列のため、他の箇条書きと同じ
        # 全角ダッシュ「—」だと項目テキストに対して不自然に長く見える。
        # 中黒「・」は幅が狭く、狭い列でも項目とのバランスが崩れない。
        add_item_list(slide, region.x, region.y + Inches(0.62), region.w,
                      region.h - Inches(0.65), phase.get("items", []),
                      bullet="・", body_gap=5)

    gates = spec.get("gates", [])
    MarkerOverlay(slide, gate_row, gates, track_y=Inches(0.4),
                 marker_color=PALETTE.blue, label_color=PALETTE.text_primary)
    _conclude(slide, conclusion, spec, default_style="editorial")


def render_table_with_conclusion(builder, spec, page):
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "dense"), page=page)
    table_region, conclusion = area.rows([4.35, 0.72], gap=Inches(0.24))
    add_data_table(slide, table_region, spec.get("columns", []),
                   spec.get("rows", []))
    _conclude(slide, conclusion, spec)


def _render_chart_visual(builder, slide, spec, region):
    """imageが指定されればPNGを、chartが指定されればネイティブグラフを描く。"""
    chart_spec = spec.get("chart")
    if chart_spec:
        add_native_chart(
            slide, region.x, region.y, region.w, region.h,
            chart_type=chart_spec.get("type", "column"),
            categories=chart_spec.get("categories", []),
            series=chart_spec.get("series", []),
            typography=_type_for(slide),
            value_format=chart_spec.get("value_format"))
    else:
        image_path = builder.paths.input_dir / spec["image"]
        add_image_contain(slide, image_path, region)


def _visual_with_insight(slide, area, spec, render_visual):
    """「主要な視覚要素＋読み取れることの箇条書き＋結論」という、
    chart_with_insight(standard)とtable_with_insightで共通する骨格。
    視覚要素そのものの描画だけをrender_visual(region)へ委譲することで、
    グラフ・表のどちらでも同じ構造を再利用する。
    """
    visual_and_notes, conclusion = area.rows([4.35, 0.72], gap=Inches(0.24))
    visual, notes = visual_and_notes.columns([1.8, 0.8], gap="wide")
    render_visual(visual)
    add_section_lead(slide, notes.x, notes.y, notes.w,
                     spec.get("insight_heading", "読み取れること"))
    add_item_list(slide, notes.x, notes.y + Inches(0.65), notes.w,
                  notes.h - Inches(0.7), spec.get("insights", []), bullet="—")
    _conclude(slide, conclusion, spec)


def render_chart_with_insight(builder, spec, page):
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    variant = spec.get("variant", "standard")
    if variant == "conclusion-led":
        insight, chart = area.columns([0.82, 1.7], gap="wide")
        add_background_zone(slide, insight.x, insight.y, insight.w, insight.h,
                            tone="brand-soft", rounded=True)
        inner = insight.inset(Inches(0.3), Inches(0.34))
        add_textbox(slide, inner.x, inner.y, inner.w, Inches(0.3),
                    spec.get("insight_heading", "読み取れること"),
                    size=_type_for(slide).small, color=PALETTE.blue, bold=True)
        add_textbox(slide, inner.x, inner.y + Inches(0.58), inner.w,
                    Inches(1.55), spec["primary_message"],
                    size=_type_for(slide).section, color=PALETTE.text_primary,
                    bold=True, line_spacing=1.02)
        if spec.get("insights"):
            add_item_list(slide, inner.x, inner.y + Inches(2.35), inner.w,
                          inner.h - Inches(2.4), spec["insights"], bullet="—")
        _render_chart_visual(builder, slide, spec, chart)
    else:
        _visual_with_insight(
            slide, area, spec,
            lambda region: _render_chart_visual(builder, slide, spec, region))


def render_table_with_insight(builder, spec, page):
    """table_with_conclusionが「結論は1文に絞る」のに対し、表から複数の
    気づきを箇条書きで示したい場合に使う（chart_with_insightの表版）。
    骨格は_visual_with_insightをそのまま共有し、視覚要素だけadd_data_table
    に差し替える。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "dense"), page=page)
    _visual_with_insight(
        slide, area, spec,
        lambda region: add_data_table(
            slide, region, spec.get("columns", []), spec.get("rows", [])))


def render_org_layers(builder, spec, page):
    """意思決定・運営など縦の責任階層と、横に並ぶ実行部門を分けて示す。

    layout-patterns.md「組織: 階層と横の役割」に対応する。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    layers = spec.get("layers", [])
    execution = spec.get("execution", [])
    weights = [1.3] * len(layers) + [2.1, 0.62]
    *layer_rows, execution_row, conclusion = area.rows(weights, gap=Inches(0.16))
    tones = ["brand-soft", "neutral", "teal-soft"]
    if layer_rows:
        layers_region = Region(
            area.x, layer_rows[0].y, area.w,
            layer_rows[-1].y + layer_rows[-1].h - layer_rows[0].y)
        layer_contents = BoxGrid(slide, layers_region, layers,
                                 rows=len(layers), skin="zone", tones=tones,
                                 row_weights=[1.3] * len(layers), gap=Inches(0.16),
                                 inset_x=Inches(0.3), inset_y=Inches(0.14))
        for layer, inner in zip(layers, layer_contents):
            # 見出し分のオフセットを固定0.46inのままにすると、layers数が増えて
            # 行の高さが縮んだ際に本文用の残りスペースをほぼ食い潰してしまう。
            # 行の高さに対する上限付き割合で縮めるが、add_section_leadの縦棒
            # マーカー自体も同じ値で縮め、offset = marker_h + GAPで揃えることで
            # マーカーが本文へ被らないことを構造的に保証する（マーカーだけ
            # 固定のまま offset だけ縮めると、マーカーが本文へ被る）。
            marker_h = min(Inches(0.38), max(Inches(0.22), int(inner.h * 0.34)))
            # 下限(0.22in)がinner.hを上回ると本文の高さが負になり、
            # python-pptxがValueErrorで落ちる（layers 7件以上で発生していた）。
            # マーカーと余白の合計が行の高さを超えないよう上限も掛ける。
            # marker_hとheader_offsetは同じ値から導くため、切り詰めても
            # マーカーが本文へ被らない関係は保たれる。
            marker_h = min(marker_h, max(0, inner.h - SECTION_LEAD_GAP))
            add_section_lead(slide, inner.x, inner.y, inner.w,
                             layer.get("heading", ""), marker_h=marker_h)
            header_offset = marker_h + SECTION_LEAD_GAP
            add_paragraph_textbox(slide, inner.x, inner.y + header_offset,
                                  inner.w, max(0, inner.h - header_offset), [
                {"segments": [(layer.get("title", ""), {
                    "size": typography.body, "color": PALETTE.text_primary,
                    "bold": True, "font": typography.body_font,
                })], "space_after": 3},
                {"segments": [(layer.get("body", ""), {
                    "size": typography.small, "color": PALETTE.text_secondary,
                    "font": typography.body_font,
                })]},
            ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    add_background_zone(slide, execution_row.x, execution_row.y,
                        execution_row.w, execution_row.h,
                        tone=tones[len(layer_rows) % len(tones)], rounded=True)
    execution_inner = execution_row.inset(Inches(0.3), Inches(0.16))
    add_section_lead(slide, execution_inner.x, execution_inner.y,
                     execution_inner.w,
                     spec.get("execution_heading", "業務実行"),
                     color=PALETTE.accent_secondary)
    cards_row = execution_inner.inset(top=Inches(0.48), bottom=0)
    execution_contents = BoxGrid(slide, cards_row, execution,
                                 inset_x=Inches(0.16), inset_y=Inches(0.16))
    for item, inner in zip(execution, execution_contents):
        add_paragraph_textbox(slide, inner.x, inner.y, inner.w, inner.h, [
            {"segments": [(item.get("title", ""), {
                "size": typography.body, "color": PALETTE.text_primary,
                "bold": True, "font": typography.body_font,
            })], "space_after": 5},
            {"segments": [(item.get("body", ""), {
                "size": typography.small, "color": PALETTE.text_secondary,
                "font": typography.body_font,
            })]},
        ], vertical_anchor=MSO_ANCHOR.MIDDLE)
    _conclude(slide, conclusion, spec)


def render_priority_actions(builder, spec, page):
    """優先度付きの課題と、対応する方針を左右に並べる。

    layout-patterns.md「課題と対応策」に対応する。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    body, conclusion = area.rows([4.35, 0.72], gap=Inches(0.24))
    left, right = body.columns([1.05, 1], gap="wide")

    issues = spec.get("issues", [])
    actions = spec.get("actions", [])

    add_section_lead(slide, left.x, left.y, left.w,
                     spec.get("issues_heading", "想定される課題"))
    top_priority = spec.get("top_priority", "最優先")
    gap = _adaptive_gap(issues, left.h - HEADING_BLOCK_H, base_gap=12,
                        per_item_pt=typography.body.pt * 1.2 +
                        typography.small.pt * 1.2 + 2)
    paragraphs = []
    for issue in issues:
        color = (PALETTE.blue if issue.get("priority") == top_priority
                 else PALETTE.text_secondary)
        paragraphs.append({"segments": [
            (f"{issue.get('priority', '')}  ", {
                "size": typography.small, "color": color, "bold": True,
                "font": typography.body_font,
            }),
            (issue.get("title", ""), {
                "size": typography.body, "color": PALETTE.text_primary,
                "bold": True, "font": typography.body_font,
            }),
        ], "space_after": 2})
        paragraphs.append({"segments": [
            (f"      {issue.get('body', '')}", {
                "size": typography.small, "color": PALETTE.text_secondary,
                "font": typography.body_font,
            }),
        ], "space_after": gap})
    add_paragraph_textbox(slide, left.x, left.y + HEADING_BLOCK_H, left.w,
                          left.h - HEADING_BLOCK_H - Inches(0.1), paragraphs)

    add_emphasis_zone(slide, right, tone="neutral")
    _lead_list(slide, right, spec.get("actions_heading", "対応方針"), actions,
              color=PALETTE.accent_secondary, bullet="—", bottom_pad=Inches(0.1))

    _conclude(slide, conclusion, spec)


def render_stage_track(builder, spec, page):
    """左から右への段階的な進行を、同格のステージカードで示す。

    段階数が2以上ならCardの間を矢印でつなぎ、順番であることを明示する
    （connectors=Falseで矢印なしのCard群に戻せる）。

    各ステージの中身はbody（1つの文章）とitems（箇条書き）のどちらでも
    書ける。両方指定した場合はbodyが上、itemsがその下に並ぶ。段階ごとに
    「やること」を複数並べる計画表はitemsを使う（文章に詰め込むより
    読み取りやすく、後からPowerPoint上で1行だけ足す編集もしやすい）。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    stages_row, note_row, conclusion = area.rows(
        [3.9, 0.34, 0.62], gap=Inches(0.16))
    stages = spec.get("stages", [])
    tones = ["neutral", "brand-soft", "teal-soft"]
    inset_x = Inches(0.28)
    contents = BoxGrid(slide, stages_row, stages, skin="zone", tones=tones,
                       gap="wide", inset_x=inset_x, inset_y=Inches(0.26))
    if spec.get("connectors", True):
        arrow_y = stages_row.y + stages_row.h // 2
        for left, right in zip(contents, contents[1:]):
            Connector(slide, left.x + left.w + inset_x, arrow_y,
                     right.x - inset_x, arrow_y, color=PALETTE.grey_500)
    blocks = []
    for index, (stage, inner) in enumerate(zip(stages, contents)):
        items = stage.get("items", [])
        width_pt = inner.w / 12700
        head_pt = 0.0
        paragraphs = []
        for text, size, color, bold, font, space_after in (
                (stage.get("label", f"STEP {index + 1}"), typography.small,
                 PALETTE.blue, True, typography.body_font, 6),
                (stage.get("title", ""), typography.section,
                 PALETTE.text_primary, True, typography.headline_font, 8)):
            paragraphs.append({"segments": [(text, {
                "size": size, "color": color, "bold": bold, "font": font,
            })], "space_after": space_after})
            head_pt += estimate_paragraph_height_pt(
                text, size.pt, width_pt, space_after=space_after)
        if stage.get("body") or not items:
            body_text = stage.get("body", "")
            paragraphs.append({"segments": [(body_text, {
                "size": typography.small, "color": PALETTE.text_secondary,
                "font": typography.body_font,
            })], "line_spacing": 1.15})
            head_pt += estimate_paragraph_height_pt(
                body_text, typography.small.pt, width_pt, line_spacing=1.15,
                space_after=6)
        list_pt = (estimate_item_list_height_pt(typography, items, width_pt,
                                                body_gap=4, title_prefix="・  ")
                   if items else 0.0)
        blocks.append({"items": items, "paragraphs": paragraphs,
                       "head_pt": head_pt, "list_pt": list_pt})

    # 箇条書きは記号の折り返し字下げをadd_item_listへ任せる（自前で段落を
    # 組むと記号の幅ぶんの字下げが失われ、2行目が記号に埋もれる）。そのぶん
    # 見出しブロックとは別のtextboxになるため、両方の高さを見積もって
    # 1つの塊として置く。塊の開始位置は全ステージ共通にし、高さは最も高い
    # ステージに合わせる。カードごとに中央へ置くと、本文の行数や項目数の
    # 違いでSTEPのラベルが数ミリずつズレて並び、横並びの同格Cardとして
    # 揃って見えない（bodyだけを使うページでも同じズレが出る）。
    if not blocks:
        top = None
    else:
        block_pt = max(block["head_pt"] + block["list_pt"] for block in blocks)
        block_h = min(contents[0].h, Inches(block_pt / 72))
        top = contents[0].y + max(0, (contents[0].h - block_h) // 2)
    for block, inner in zip(blocks, contents):
        bottom = inner.y + inner.h
        head_h = min(bottom - top, Inches(block["head_pt"] / 72))
        add_paragraph_textbox(slide, inner.x, top, inner.w, head_h,
                              block["paragraphs"])
        if block["items"]:
            list_h = min(bottom - top - head_h, Inches(block["list_pt"] / 72))
            add_item_list(slide, inner.x, top + head_h, inner.w, list_h,
                          block["items"], bullet="・", body_gap=4)
    if spec.get("note"):
        add_textbox(slide, note_row.x, note_row.y, note_row.w, note_row.h,
                    spec["note"], size=typography.small,
                    color=PALETTE.text_secondary)
    _conclude(slide, conclusion, spec)


def render_numbered_list(builder, spec, page):
    """番号付きの項目群を、アジェンダや依頼事項のように単列で示す。

    項目数が少なくても余白が偏らないよう、リスト全体を利用可能な高さの
    中央へ配置する。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    items = spec.get("items", [])
    position = spec.get("message_position", "top")
    default_style = "plain" if position == "top" else "solid"
    style = spec.get("message_style", default_style)

    if position == "top":
        add_key_message(slide, area.x, area.y, area.w,
                        spec["primary_message"], style=style)
        list_top = area.y + Inches(0.62)
        list_bottom = area.y + area.h
    else:
        list_top = area.y
        list_bottom = area.y + area.h - Inches(1.1)

    available = max(0, list_bottom - list_top)
    # 項目が少なくても余白が偏らないよう、add_text_list側でブロックを
    # 利用可能な高さの中央へ配置する（項目数が多い場合はmax_row_hを
    # 上限に行の高さを縮め、スライド外へのはみ出しを防ぐ）。
    add_text_list(slide, area.x, list_top, area.w, available, items,
                 marker="number", divider=True, max_row_h=Inches(0.9))

    if position == "bottom":
        add_key_message(slide, area.x, list_bottom + Inches(0.25), area.w,
                        spec["primary_message"], style=style)


def render_section_divider(builder, spec, page):
    builder.add_section_divider(spec["title"], kicker=spec.get("kicker"),
                                subtitle=spec.get("subtitle"), page=page)


def render_matrix(builder, spec, page):
    """rows×colsのマス目に分け、各マスの位置づけを示す(ポートフォリオ分析等)。

    x_axis/y_axisを指定すると連続軸（低い⇄高いの度合い）上の位置づけ
    として使う（BCGの成長率-シェア、Eisenhowerの緊急度-重要度等）。
    どちらも指定しなければ、SWOT等「軸を持たない固定カテゴリ」として
    軸ラベル分の余白を使わずマス目を広く使う。専用のaxesフラグは持たず
    x_axis/y_axisの有無だけで自動的に切り替わる（値を与えれば使う、
    という以上の情報をこの2つは持たないため）。

    rows/colsは省略時2x2（BCGマトリクスやSWOT等、最も一般的な2軸/
    4象限フレームワークに合わせた既定）。GE-McKinseyの3x3のような
    2x2を超えるマトリクスも、同じ構造でrows/colsを指定するだけで
    作れる（次元数はBoxGridと同じく1つの操作のパラメータに過ぎない）。

    cellsの並び順は左上から右へ、次の行へ、という読み順。先頭の行が
    y_axis.high側（上）、左の列がx_axis.low側になる。上下を取り違えると
    象限の意味が反転するため注意する（renderer-catalog.mdにも明記）。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    rows = spec.get("rows", 2)
    cols = spec.get("cols", 2)
    x_axis = spec.get("x_axis") or {}
    y_axis = spec.get("y_axis") or {}
    axes = bool(x_axis) or bool(y_axis)
    if axes:
        grid_area, caption_row, conclusion = area.rows(
            [4.5, 0.34, 0.62], gap=Inches(0.16))
        axis_col, plot_col = grid_area.columns([0.14, 1], gap=Inches(0.14))
        add_textbox(slide, axis_col.x, axis_col.y, axis_col.w, Inches(0.4),
                    y_axis.get("high", ""), size=typography.small,
                    color=PALETTE.text_secondary, bold=True)
        add_textbox(slide, axis_col.x, axis_col.y + axis_col.h - Inches(0.4),
                    axis_col.w, Inches(0.4), y_axis.get("low", ""),
                    size=typography.small, color=PALETTE.text_secondary, bold=True)
    else:
        plot_col, conclusion = area.rows([4.84, 0.62], gap=Inches(0.16))

    cells = spec.get("cells", [])
    emphasis_tone = (lambda item, i:
                     "brand-soft" if isinstance(item, dict) and item.get("emphasis")
                     else "neutral")
    contents = BoxGrid(slide, plot_col, cells, rows=rows, cols=cols,
                       skin="zone", tones=emphasis_tone, gap=Inches(0.1),
                       inset_x=Inches(0.24), inset_y=Inches(0.2))
    for cell, inner in zip(cells, contents):
        add_paragraph_textbox(slide, inner.x, inner.y, inner.w, inner.h, [
            {"segments": [(cell.get("label", ""), {
                "size": typography.small, "color": PALETTE.blue,
                "bold": True, "font": typography.body_font,
            })], "space_after": 4},
            {"segments": [(cell.get("title", ""), {
                "size": typography.body, "color": PALETTE.text_primary,
                "bold": True, "font": typography.body_font,
            })], "space_after": 3},
            {"segments": [(cell.get("body", ""), {
                "size": typography.small, "color": PALETTE.text_secondary,
                "font": typography.body_font,
            })]},
        ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    if axes:
        half = plot_col.w // 2
        add_textbox(slide, plot_col.x, caption_row.y, half, caption_row.h,
                    x_axis.get("low", ""), size=typography.small,
                    color=PALETTE.text_secondary, bold=True)
        add_textbox(slide, plot_col.x + half, caption_row.y, half, caption_row.h,
                    x_axis.get("high", ""), size=typography.small,
                    color=PALETTE.text_secondary, bold=True, align=PP_ALIGN.RIGHT)

    _conclude(slide, conclusion, spec)


def render_stat_highlight(builder, spec, page):
    """単一の実績数値を主役にする（stat指定時）か、複数指標を均等な
    グリッドで一覧するKPIダッシュボード（stat省略時）を描く。

    どちらも同じBoxGrid＋Statの組み合わせで、hero（主役として大きく
    見せる1指標）の有無だけが違う。パラメータ差で別rendererに分けない
    という方針（BandStack/BoxGridの統合と同じ理由）で、1つのrenderer
    が両方を担う。

    stat/supportingの各項目に"tone": "positive"/"negative"/"warning"を
    指定すると、数値の色が符号専用のセマンティックカラーになる（数値
    文字列の先頭が"-"かどうかでは自動判定しない。削減率のように、
    マイナスの数値が良い結果を意味することがあるため、toneは呼び出し側
    が意味を判断して明示する）。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    stat = spec.get("stat")
    supporting = spec.get("supporting", [])

    if stat:
        if supporting:
            hero_row, supporting_row, conclusion = area.rows(
                [2.5, 2.05, 0.62], gap=Inches(0.2))
        else:
            hero_row, conclusion = area.rows([4.7, 0.62], gap=Inches(0.24))
            supporting_row = None
        add_background_zone(slide, hero_row.x, hero_row.y, hero_row.w, hero_row.h,
                            tone="brand-soft", rounded=True)
        inner = hero_row.inset(Inches(0.5), Inches(0.3))
        Stat(slide, inner.x, inner.y, inner.w, inner.h,
            stat.get("value", ""), stat.get("label", ""), detail=stat.get("detail"),
            value_size=Pt(56), value_color=_tone_color(stat, PALETTE.navy),
            label_size=typography.section, vertical_anchor=MSO_ANCHOR.MIDDLE)
    else:
        supporting_row, conclusion = area.rows([4.7, 0.62], gap=Inches(0.24))

    if supporting_row is not None and supporting:
        cols = min(4, len(supporting))
        value_size = typography.metric if stat else Pt(34)
        contents = BoxGrid(slide, supporting_row, supporting, cols=cols,
                           inset_x=Inches(0.2), inset_y=Inches(0.2))
        for item, inner in zip(supporting, contents):
            Stat(slide, inner.x, inner.y, inner.w, inner.h,
                item.get("value", ""), item.get("label", ""),
                value_size=value_size, value_color=_tone_color(item, PALETTE.blue),
                label_size=typography.small, label_color=PALETTE.text_secondary,
                label_bold=False, vertical_anchor=MSO_ANCHOR.MIDDLE)

    _conclude(slide, conclusion, spec)


def render_funnel(builder, spec, page):
    """順を追って絞り込まれていく推移を、段ごとの値に応じた帯の幅で示す
    （リード獲得のファネル分析、市場規模のTAM/SAM/SOM等）。

    stagesは値の大きい順（絞り込みが進むほど値が小さくなる順）に並べる。
    帯の幅は正確な比率ではなくおおよその絞り込み具合を示す構造表現の
    ため、比率そのものを厳密に伝えたい場合はchart_with_insightの
    棒グラフを使う。

    insightsを指定すると、帯を_visual_with_insightの視覚要素として
    右側に気づきの箇条書きを添える（table_with_insight/chart_with_insight
    と同じ骨格）。段の値だけでなく、そこから読み取れる複数の気づきを
    合わせて示したい場合はこちらを使う。省略時は帯が全幅を使う。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    stages = spec.get("stages", [])
    tones = ["brand-soft", "teal-soft", "neutral", "brand-soft"]

    def draw_stack(region):
        contents = ProportionalStack(slide, region, stages, skin="zone",
                                     tones=tones, gap=Inches(0.12))
        for stage, inner in zip(stages, contents):
            add_paragraph_textbox(slide, inner.x, inner.y, inner.w, inner.h, [
                {"segments": [(stage.get("title", ""), {
                    "size": typography.body, "color": PALETTE.text_primary,
                    "bold": True, "font": typography.body_font,
                })], "space_after": 3, "align": PP_ALIGN.CENTER},
                {"segments": [(stage.get("value_label", ""), {
                    "size": typography.small, "color": PALETTE.text_secondary,
                    "font": typography.body_font,
                })], "align": PP_ALIGN.CENTER},
            ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    if spec.get("insights"):
        _visual_with_insight(slide, area, spec, draw_stack)
    else:
        stack_row, conclusion = area.rows([4.7, 0.62], gap=Inches(0.22))
        draw_stack(stack_row)
        _conclude(slide, conclusion, spec)


def _cycle_box_w(steps, square, typography):
    """cycleのCard幅を、最も長いtitleが1行に収まる幅へ広げる。

    固定幅のままだと「ガイドライン改善」が「ガイドライン改／善」のように
    語の途中で折り返す。一方で広げすぎると円周上の隣接Card同士が接触する
    ため、隣り合う中心間の距離から上限を決める（半径は概ね外接正方形の
    1/4なので、中心間距離は 2*r*sin(pi/n)）。
    """
    if not steps:
        return Inches(1.5)
    longest = max((str(step.get("title", "")) for step in steps), key=len,
                  default="")
    # 最も長いtitleが1行に収まる幅（pt）＋左右のinsetぶん。実際に置ける
    # 幅の上限はRadialLayoutが配置から決めるため、ここでは要望値だけ返す。
    needed_pt = sum(char_width_factor(ch) for ch in longest) * typography.body.pt
    needed = Inches(needed_pt / 72) + Inches(0.36)
    # 上限は正方形の半分。これを超えると円というより横並びに見えるため。
    return int(max(Inches(1.5), min(needed, square.w // 2)))


def render_cycle(builder, spec, page):
    """繰り返し・循環するプロセス（PDCA等）を、円周上に並べた同格のCard群
    と、隣接する項目を結ぶ矢印で示す（最後尾から先頭へも矢印で結び輪に
    する）。stage_trackの一方向の進行とは異なり、繰り返しであることが
    要点の場合に使う。stepsは4〜6件程度を目安にする（多いと隣接する
    Card同士の間隔が狭くなり、矢印やテキストが読みにくくなる）。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    circle_row, conclusion = area.rows([4.7, 0.62], gap=Inches(0.24))
    size = min(circle_row.w, circle_row.h)
    square = Region(circle_row.x + (circle_row.w - size) // 2,
                    circle_row.y + (circle_row.h - size) // 2, size, size)
    steps = spec.get("steps", [])
    tones = ["brand-soft", "teal-soft", "neutral", "brand-soft", "teal-soft", "neutral"]
    contents = RadialLayout(slide, square, steps, tones=tones,
                            box_w=_cycle_box_w(steps, square, typography))
    for index, (step, inner) in enumerate(zip(steps, contents)):
        add_paragraph_textbox(slide, inner.x, inner.y, inner.w, inner.h, [
            {"segments": [(step.get("label", f"STEP {index + 1}"), {
                "size": typography.small, "color": PALETTE.blue,
                "bold": True, "font": typography.body_font,
            })], "space_after": 2, "align": PP_ALIGN.CENTER},
            {"segments": [(step.get("title", ""), {
                "size": typography.body, "color": PALETTE.text_primary,
                "bold": True, "font": typography.body_font,
            })], "align": PP_ALIGN.CENTER},
        ], vertical_anchor=MSO_ANCHOR.MIDDLE)
    _conclude(slide, conclusion, spec)


def render_timeline(builder, spec, page):
    """期間を持つ複数の取り組みを、共通の時間軸上の帯で示す
    （ロードマップ、実行計画のガント）。

    帯の「長さ」が期間を表す点が他typeとの違い。stage_trackは段階の
    順序だけを、process_with_gatesは判断を下す時点だけを示し、どちらも
    長さに意味を持たせない。開始と終了が異なる取り組みが並走する計画は
    このtypeを使う。判断ポイントを合わせて示したい場合は
    process_with_gatesのページを別に立てる（1ページに詰め込まない）。

    periodsは時間軸の目盛り（["Q1", "Q2", ...]）。rowsの各要素は
    start/end（1始まりのperiod番号、endを含む）でその区間を占める。
    endを省略するとstartと同じ1期間。範囲外の値はperiodsの範囲へ丸める
    （生成を止めず、目に見える形で端に寄せる）。

    帯に収まらない長さのtitleは帯の右側へ回す。狭い帯へ押し込むと文字が
    溢れて読めなくなるため、帯の高さに入る行数を超える場合は外へ出す。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    periods = spec.get("periods", [])
    rows = spec.get("rows", [])
    header_row, lane_area, conclusion = area.rows(
        [0.42, 4.2, 0.62], gap=Inches(0.16))
    # 見出し行とレーン行を同じ重み・同じgapで分割すると、両者のx座標が
    # 一致する（Regionの分割は決定的なため）。目盛りと帯がズレない。
    split = [1.5, 6.0]
    _, header_track = header_row.columns(split, gap=Inches(0.14))
    label_col, track_col = lane_area.columns(split, gap=Inches(0.14))

    count = max(1, len(periods))
    unit = track_col.w / count
    add_background_zone(slide, header_track.x, header_row.y, header_track.w,
                        header_row.h, tone="neutral", rounded=True)
    for index, period in enumerate(periods):
        add_textbox(slide, int(track_col.x + unit * index),
                    header_row.y + Inches(0.06), int(unit), header_row.h,
                    str(period), size=typography.small,
                    color=PALETTE.text_primary, bold=True, align=PP_ALIGN.CENTER)
        if index:
            # 期間の区切り線。add_hairlineは横罫線専用のため、縦は細い
            # Boxで引く（Connectorは関係を示す線なので区切りには使わない）。
            Box(slide, int(track_col.x + unit * index), lane_area.y,
                Pt(0.75), lane_area.h, rounded=False, fill=PALETTE.line_neutral)

    lane_weights = [1] * max(1, len(rows))
    label_lanes = label_col.rows(lane_weights, gap=Inches(0.14))
    track_lanes = track_col.rows(lane_weights, gap=Inches(0.14))
    default_tones = ["brand", "teal", "neutral"]
    track_right = track_col.x + track_col.w
    for index, (row, label_lane, lane) in enumerate(
            zip(rows, label_lanes, track_lanes)):
        add_paragraph_textbox(slide, label_lane.x, label_lane.y, label_lane.w,
                              label_lane.h, [
            {"segments": [(row.get("label", ""), {
                "size": typography.body, "color": PALETTE.text_primary,
                "bold": True, "font": typography.body_font,
            })]},
        ], vertical_anchor=MSO_ANCHOR.MIDDLE)

        start = _clamp_period(row.get("start", 1), count)
        end = max(start, _clamp_period(row.get("end", start), count))
        pad = Inches(0.05)
        bar_x = int(track_col.x + unit * (start - 1)) + pad
        bar_w = max(Inches(0.12), int(unit * (end - start + 1)) - 2 * pad)
        # レーンが少ないページで帯が間延びしないよう高さに上限を設け、
        # レーンの中央へ置く。
        bar_h = min(lane.h, Inches(0.62))
        bar_y = lane.y + (lane.h - bar_h) // 2
        tone = row.get("tone", default_tones[index % len(default_tones)])
        add_background_zone(slide, bar_x, bar_y, bar_w, bar_h,
                            tone=_ZONE_TONES.get(tone, "brand-soft"), rounded=True)

        title = row.get("title", "")
        text_pad = Inches(0.14)
        inner_w_pt = max(1, (bar_w - 2 * text_pad) / 12700)
        line_h_pt = typography.body.pt * 1.25
        max_lines = max(1, int((bar_h / 12700) / line_h_pt))
        overflows = estimate_line_count(title, typography.body.pt, inner_w_pt) > max_lines
        outside_w = track_right - (bar_x + bar_w) - Inches(0.12)
        if overflows and outside_w > Inches(0.7):
            text_x, text_w = bar_x + bar_w + Inches(0.12), outside_w
        else:
            text_x, text_w = bar_x + text_pad, bar_w - 2 * text_pad
        add_paragraph_textbox(slide, text_x, bar_y, text_w, bar_h, [
            {"segments": [(title, {
                "size": typography.body, "color": PALETTE.text_primary,
                "bold": True, "font": typography.body_font,
            })], "line_spacing": 1.05},
        ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    _conclude(slide, conclusion, spec)


def render_waterfall(builder, spec, page):
    """開始値から増減の要因を積み上げ、終了値へどう到達したかを示す
    （ブリッジ図／ウォーターフォール。売上・利益の増減要因分解、予算差異、
    価格/数量/構成の分解）。

    barsは左から右へ順に読む。kind="total"の棒は基準線0からの絶対値
    （開始値・終了値・小計）として描き、kindを省略した棒は直前までの
    累計に対する増減として宙に浮かせる。増減の色は既定ではvalueの符号で
    決まる（増=positive、減=negative）。棒と棒の間には、累計の水準を示す
    細い連結線を引く。

    ただしコストのブリッジのように「減ることが良い結果」の場合は符号と
    意味が逆になるため、棒ごとに"tone": "positive"/"negative"/"warning"
    で色を明示できる（stat_highlightのtoneと同じ語彙・同じ考え方。数値の
    符号だけで良し悪しを決めない）。

    「合計の内訳」を並べたいだけならchart_with_insightの積み上げ棒を
    使う。このtypeはAからBへ変化した理由を分解して示す図。

    value_labelを指定しない場合はvalueから自動生成する（増減には符号を
    付ける）。単位付きの表記（「+1.8億」等）にしたい場合は指定する。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    bars = spec.get("bars", [])
    plot_row, label_row, conclusion = area.rows(
        [3.75, 0.55, 0.62], gap=Inches(0.18))

    segments, cumulative = [], 0.0
    for bar in bars:
        value = _number(bar.get("value", 0))
        if bar.get("kind") == "total":
            bottom, top, cumulative, kind = min(0.0, value), max(0.0, value), value, "total"
        else:
            start, end = cumulative, cumulative + value
            bottom, top = min(start, end), max(start, end)
            cumulative = end
            kind = "increase" if value >= 0 else "decrease"
        segments.append({"bar": bar, "value": value, "kind": kind,
                         "bottom": bottom, "top": top, "level": cumulative})

    low = min([segment["bottom"] for segment in segments] + [0.0])
    high = max([segment["top"] for segment in segments] + [0.0])
    # 全ての値が0のとき（YAMLの書きかけ等）でも0除算で生成を止めない。
    # 棒は高さ0となり、値が入っていないことがページ上で見て分かる。
    span = (high - low) or 1.0
    # 値ラベルは棒の上に置くため、最上段の棒の上へその分の帯を空ける。
    bars_area = plot_row.inset(0, Inches(0.32), 0, 0)

    def to_y(value):
        return int(bars_area.y + bars_area.h * (high - value) / span)

    if low < 0:
        add_hairline(slide, bars_area.x, to_y(0), bars_area.w,
                     color=PALETTE.grey_300)

    slots = bars_area.columns([1] * max(1, len(bars)), gap="tight")
    label_slots = label_row.columns([1] * max(1, len(bars)), gap="tight")
    colors = {"total": PALETTE.navy, "increase": PALETTE.positive,
              "decrease": PALETTE.negative}
    previous = None
    for segment, slot, label_slot in zip(segments, slots, label_slots):
        top_y, bottom_y = to_y(segment["top"]), to_y(segment["bottom"])
        bar_w = int(slot.w * 0.66)
        bar_x = slot.x + (slot.w - bar_w) // 2
        color = _tone_color(segment["bar"], colors[segment["kind"]])
        if previous is not None:
            level_y = to_y(previous["level"])
            add_hairline(slide, previous["right"], level_y,
                         max(0, bar_x - previous["right"]), color=PALETTE.grey_300)
        # 値が0の棒も存在自体は見えるようにする（高さ0だと棒が消え、
        # 項目を書いたのに描かれていないように見えるため）。
        Box(slide, bar_x, top_y, bar_w, max(Pt(1.5), bottom_y - top_y),
            rounded=False, fill=color)
        add_textbox(slide, slot.x, top_y - Inches(0.3), slot.w, Inches(0.28),
                    segment["bar"].get("value_label")
                    or _delta_label(segment["value"], segment["kind"]),
                    size=typography.small, color=color, bold=True,
                    align=PP_ALIGN.CENTER)
        is_total = segment["kind"] == "total"
        add_paragraph_textbox(slide, label_slot.x, label_slot.y, label_slot.w,
                              label_slot.h, [
            {"segments": [(segment["bar"].get("label", ""), {
                "size": typography.small,
                "color": PALETTE.text_primary if is_total else PALETTE.text_secondary,
                "bold": is_total, "font": typography.body_font,
            })], "align": PP_ALIGN.CENTER, "line_spacing": 1.05},
        ])
        previous = {"level": segment["level"], "right": bar_x + bar_w}

    _conclude(slide, conclusion, spec)


def render_issue_tree(builder, spec, page):
    """論点を階層的に分解して示す（イシューツリー／ロジックツリー）。

    左に根（分解する問い）、中央に第1階層の分解軸、右にその内訳を置く。
    階層は3段までに固定する。それ以上は1ページの幅に収まらず、読み手も
    枝を追えなくなるため、深い分解は章を分けて複数ページにする。

    branchesがitemsを1つも持たない場合は2段のツリーとして描き、右の列を
    使わずに中央の枝を広く取る（空の列を残すと余白が偏るため）。

    線は矢印を付けない。ツリーは時間や因果の流れではなく「分解した結果
    の包含関係」を示すため、向きを持たせない（順序を示したい場合は
    stage_trackやprocess_with_gatesを使う）。
    """
    slide, area = builder.add_slide(
        spec["title"], density=spec.get("density", "standard"), page=page)
    typography = _type_for(slide)
    root = spec.get("root") or {}
    branches = spec.get("branches", [])
    body, conclusion = area.rows([4.7, 0.62], gap=Inches(0.24))
    has_leaves = any(branch.get("items") for branch in branches)
    columns = body.columns([1.15, 1.3, 1.55] if has_leaves else [1.15, 2.2],
                           gap="wide")
    root_col, branch_col = columns[0], columns[1]
    leaf_col = columns[2] if has_leaves else None

    root_h = min(body.h, Inches(1.95))
    root_y = body.y + (body.h - root_h) // 2
    add_background_zone(slide, root_col.x, root_y, root_col.w, root_h,
                        tone="brand-soft", rounded=True)
    root_inner = Region(root_col.x, root_y, root_col.w, root_h).inset(
        Inches(0.24), Inches(0.2))
    root_paragraphs = []
    if root.get("label"):
        root_paragraphs.append({"segments": [(root["label"], {
            "size": typography.small, "color": PALETTE.blue,
            "bold": True, "font": typography.body_font,
        })], "space_after": 5})
    root_paragraphs.append({"segments": [(root.get("title", ""), {
        "size": typography.section, "color": PALETTE.text_primary,
        "bold": True, "font": typography.headline_font,
    })], "space_after": 6, "line_spacing": 1.1})
    if root.get("body"):
        root_paragraphs.append({"segments": [(root["body"], {
            "size": typography.small, "color": PALETTE.text_secondary,
            "font": typography.body_font,
        })], "line_spacing": 1.15})
    add_paragraph_textbox(slide, root_inner.x, root_inner.y, root_inner.w,
                          root_inner.h, root_paragraphs,
                          vertical_anchor=MSO_ANCHOR.MIDDLE)

    # BoxGridが返すのは内側の余白を差し引いたRegionのため、Cardの外形
    # （＝線を結ぶ位置）はinsetを足し戻して求める。両者がズレると線が
    # Cardの縁から浮くので、insetは変数を共有して同じ値から導く。
    branch_inset_x, branch_inset_y = Inches(0.22), Inches(0.16)
    contents = BoxGrid(slide, branch_col, branches, rows=max(1, len(branches)),
                       cols=1, gap=Inches(0.16), inset_x=branch_inset_x,
                       inset_y=branch_inset_y)
    root_center_y = root_y + root_h // 2
    leaf_bands = None
    if leaf_col is not None:
        # 枝ごとの帯を内訳の件数に比例させる。枝は等分のままなので中心は
        # 揃わないが、線はelbowで結ぶため対応関係は保たれる。
        leaf_bands = leaf_col.rows(
            [max(1, len(branch.get("items", []))) for branch in branches],
            gap=Inches(0.18))
    for index, (branch, inner) in enumerate(zip(branches, contents)):
        center_y = (inner.y - branch_inset_y) + (inner.h + 2 * branch_inset_y) // 2
        Connector(slide, root_col.x + root_col.w, root_center_y,
                  inner.x - branch_inset_x, center_y, style="elbow",
                  arrow="none", color=PALETTE.grey_500, width=Pt(1.25))
        paragraphs = [{"segments": [(branch.get("title", ""), {
            "size": typography.body, "color": PALETTE.text_primary,
            "bold": True, "font": typography.body_font,
        })], "space_after": 4, "line_spacing": 1.1}]
        if branch.get("body"):
            paragraphs.append({"segments": [(branch["body"], {
                "size": typography.small, "color": PALETTE.text_secondary,
                "font": typography.body_font,
            })], "line_spacing": 1.15})
        add_paragraph_textbox(slide, inner.x, inner.y, inner.w, inner.h,
                              paragraphs, vertical_anchor=MSO_ANCHOR.MIDDLE)

        if leaf_bands is None:
            continue
        items = branch.get("items", [])
        band = leaf_bands[index]
        leaf_rows = band.rows([1] * max(1, len(items)), gap=Inches(0.08))
        for item, leaf in zip(items, leaf_rows):
            add_background_zone(slide, leaf.x, leaf.y, leaf.w, leaf.h,
                                tone="neutral", rounded=True)
            Connector(slide, inner.x + inner.w + branch_inset_x, center_y,
                      leaf.x, leaf.y + leaf.h // 2, style="elbow",
                      arrow="none", color=PALETTE.line_neutral, width=Pt(1))
            add_paragraph_textbox(slide, leaf.x + Inches(0.18), leaf.y,
                                  leaf.w - Inches(0.32), leaf.h, [
                {"segments": [(item if isinstance(item, str)
                               else item.get("title", ""), {
                    "size": typography.small, "color": PALETTE.text_primary,
                    "font": typography.body_font,
                })], "line_spacing": 1.05},
            ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    _conclude(slide, conclusion, spec)


RENDERERS = {
    "cover": render_cover,
    "comparison": render_comparison,
    "evidence_and_decision": render_evidence_and_decision,
    "scope_and_exclusions": render_scope_and_exclusions,
    "process_with_gates": render_process_with_gates,
    "table_with_conclusion": render_table_with_conclusion,
    "table_with_insight": render_table_with_insight,
    "chart_with_insight": render_chart_with_insight,
    "org_layers": render_org_layers,
    "priority_actions": render_priority_actions,
    "stage_track": render_stage_track,
    "numbered_list": render_numbered_list,
    "section_divider": render_section_divider,
    "matrix": render_matrix,
    "stat_highlight": render_stat_highlight,
    "funnel": render_funnel,
    "timeline": render_timeline,
    "waterfall": render_waterfall,
    "issue_tree": render_issue_tree,
    "cycle": render_cycle,
}


def render_deck(content, root, output_path=None, *, extra_renderers=None):
    """contentからデッキを生成する。

    extra_renderersは {type名: 関数(builder, spec, page)} の辞書。該当する
    rendererが無いページを生成スクリプト側で個別構築する場合に渡す
    （renderer-catalog.mdのEscape Hatch）。渡したtypeはpreflightでも
    許容され、共通の契約（title・primary_message・density・文字量・項目数）
    は他のページと同じように検査される。個別構築ページのためだけに
    preflightごと外す必要はない。

    既存typeと同じ名前は受け付けない。デッキ単位で標準rendererの意味が
    黙って差し替わると、同じtypeが資料によって別の見た目になるため。
    rendererそのものを変えたい場合はslidekit側を直す（SKILL.mdのREVISE）。
    """
    extra_renderers = extra_renderers or {}
    collisions = sorted(set(extra_renderers) & set(RENDERERS))
    if collisions:
        raise ValueError(
            "extra_renderersが既存のtypeと重複しています: "
            + ", ".join(collisions)
            + "（別の名前にするか、slidekit側のrendererを修正してください）")
    warnings = require_valid_content(content, extra_types=extra_renderers.keys())
    builder = DeckBuilder(root, mode=content.get("deck", {}).get("mode"))
    renderers = {**RENDERERS, **extra_renderers}
    for page, spec in enumerate(content["slides"], 1):
        renderers[spec["type"]](builder, spec, page)
    return builder.save(output_path), warnings
