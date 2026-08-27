#!/usr/bin/env python3
import sys
import tempfile
from pathlib import Path

import yaml
from pptx import Presentation


def _check_chart_stability(add_native_chart, TYPE_BUSINESS):
    """add_native_chartのエッジケース耐性を検証する（実際にKeynoteで

    はみ出し・重なりを確認して見つけた回帰を防ぐ）。
    """
    from pptx import Presentation as _Presentation
    from pptx.util import Inches as _Inches

    prs = _Presentation()
    prs.slide_width = _Inches(13.333333)
    prs.slide_height = _Inches(7.5)
    blank = prs.slide_layouts[6]

    def expect_error(label, **kwargs):
        slide = prs.slides.add_slide(blank)
        try:
            add_native_chart(slide, _Inches(0.5), _Inches(0.5), _Inches(6),
                             _Inches(4), typography=TYPE_BUSINESS, **kwargs)
        except ValueError:
            return
        raise AssertionError(f"{label}: ValueErrorが発生しませんでした")

    expect_error("categories/values数の不一致", chart_type="column",
                categories=["A", "B", "C"],
                series=[{"name": "s1", "values": [1, 2]}])
    expect_error("categoriesが空", chart_type="column",
                categories=[], series=[{"name": "s1", "values": []}])
    expect_error("pieに複数系列", chart_type="pie",
                categories=["A", "B"],
                series=[{"name": "s1", "values": [1, 2]}, {"name": "s2", "values": [3, 4]}])

    # 負の値: ラベルが軸のカテゴリ名と重ならないよう、棒の内側に置かれる
    # ことをXMLレベルで確認する（Keynote実測でOUTSIDE_ENDが重なるのを確認済み）。
    from pptx.enum.chart import XL_LABEL_POSITION
    neg_slide = prs.slides.add_slide(blank)
    frame = add_native_chart(neg_slide, _Inches(0.5), _Inches(0.5), _Inches(6),
                             _Inches(4), typography=TYPE_BUSINESS, chart_type="column",
                             categories=["A", "B"],
                             series=[{"name": "s1", "values": [-5, 10]}])
    assert frame.chart.plots[0].data_labels.position == XL_LABEL_POSITION.INSIDE_END, (
        "負の値を含む棒グラフのラベル位置がINSIDE_ENDになっていません")

    # カテゴリ数が多い折れ線: 点ごとのラベルが密集して読めなくなるのを防ぐため、
    # 閾値を超えたらラベルを出さない。
    many_slide = prs.slides.add_slide(blank)
    frame = add_native_chart(many_slide, _Inches(0.5), _Inches(0.5), _Inches(6),
                             _Inches(4), typography=TYPE_BUSINESS, chart_type="line",
                             categories=[f"C{i}" for i in range(20)],
                             series=[{"name": "s1", "values": list(range(20))}])
    assert frame.chart.plots[0].has_data_labels is False, (
        "カテゴリ数が多い折れ線グラフでデータラベルが抑制されていません")

    # stacked_column: セグメントが重なるようoverlap=100、ラベルは
    # セグメント内(CENTER)に収める。
    stacked_slide = prs.slides.add_slide(blank)
    frame = add_native_chart(stacked_slide, _Inches(0.5), _Inches(0.5), _Inches(6),
                             _Inches(4), typography=TYPE_BUSINESS,
                             chart_type="stacked_column",
                             categories=["A", "B"],
                             series=[{"name": "s1", "values": [3, 4]},
                                     {"name": "s2", "values": [2, 1]}])
    assert frame.chart.plots[0].overlap == 100, (
        "stacked_columnのoverlapが100になっていません")
    assert frame.chart.plots[0].data_labels.position == XL_LABEL_POSITION.CENTER, (
        "stacked_columnのラベル位置がCENTERになっていません")
    expect_error("stacked_columnのcategories/values数不一致",
                chart_type="stacked_column", categories=["A", "B"],
                series=[{"name": "s1", "values": [1, 2, 3]}])

    # scatter: categoriesを使わずpoints(x, y)を使うXY座標型。
    scatter_slide = prs.slides.add_slide(blank)
    frame = add_native_chart(scatter_slide, _Inches(0.5), _Inches(0.5), _Inches(6),
                             _Inches(4), typography=TYPE_BUSINESS,
                             chart_type="scatter",
                             series=[{"name": "s1", "points": [
                                 {"x": 1.0, "y": 2.0}, {"x": 2.0, "y": 3.5}]}])
    assert len(list(frame.chart.plots[0].series)) == 1, (
        "scatterの系列数が想定と異なります")
    expect_error("scatterのseriesが空", chart_type="scatter", series=[])
    expect_error("scatterのpointsが空", chart_type="scatter",
                series=[{"name": "s1", "points": []}])
    expect_error("scatterのpointsにyが無い", chart_type="scatter",
                series=[{"name": "s1", "points": [{"x": 1.0}]}])


def _check_edge_case_stability(content):
    """項目数・文字量・重みの極端な値でも、生成が例外で止まらないことを確認する。

    preflightが止めない入力は必ず描画できる、というのがこのスキルの契約。
    レイアウトは利用可能な領域から余白や見出し分を引いて寸法を決めるため、
    項目が多いと引きすぎて負になり、重みが全て0だとゼロ除算になりうる。
    どちらもデッキ全体の生成が失敗するため、退化した見た目に落とすことで
    回避している（過去にorg_layers 7件以上、stat_highlight 40件以上、
    phases/columnsのweight=0で実際に落ちていた）。
    """
    import copy
    from slidekit import DeckBuilder
    from slidekit.preflight import inspect_content
    from slidekit.renderers import RENDERERS

    base = {}
    for slide in content["slides"]:
        base.setdefault(slide["type"], slide)
    # ここで見るのはレイアウトの安定性のため、外部PNGへの依存は外して
    # ネイティブchartへ置き換える（画像の有無で結果が変わらないように）。
    if "image" in base.get("chart_with_insight", {}):
        chart_spec = copy.deepcopy(base["chart_with_insight"])
        chart_spec.pop("image")
        chart_spec["chart"] = {"type": "column", "categories": ["A", "B"],
                               "series": [{"name": "s", "values": [1, 2]}]}
        base["chart_with_insight"] = chart_spec
    long_text = "非常に長い説明文がここに入ります。" * 20
    skip_text = {"type", "id", "density", "variant", "message_position"}

    for slide_type, spec in sorted(base.items()):
        for count in (0, 1, 7, 25, 60):
            for long_form in (False, True):
                trial = copy.deepcopy(spec)
                for key, value in list(trial.items()):
                    if isinstance(value, list) and value:
                        trial[key] = [copy.deepcopy(value[i % len(value)])
                                      for i in range(count)]
                    elif isinstance(value, list):
                        trial[key] = []
                    elif long_form and isinstance(value, str) and key not in skip_text:
                        trial[key] = long_text
                errors, _ = inspect_content(
                    {"deck": content["deck"], "slides": [trial]})
                if errors:
                    continue  # preflightが止める入力は描画対象外
                for mode in ("business", "dense", "large-room"):
                    try:
                        RENDERERS[slide_type](
                            DeckBuilder(Path.cwd(), mode=mode), trial, 1)
                    except Exception as exc:  # noqa: BLE001
                        raise AssertionError(
                            f"{slide_type}: 項目{count}件 / 長文={long_form} / "
                            f"{mode} で生成が失敗しました: "
                            f"{type(exc).__name__}: {exc}") from exc

    # 重み・件数の退化した値（YAML由来で起こりうる）でも落ちないこと。
    degenerate = [
        {"type": "process_with_gates", "title": "T", "primary_message": "M",
         "density": "dense",
         "phases": [{"title": "p", "label": "P", "weight": 0, "items": ["x"]},
                    {"title": "q", "label": "Q", "weight": 0, "items": ["y"]}],
         "gates": [{"title": "g", "position": 0.5}]},
        {"type": "table_with_conclusion", "title": "T", "primary_message": "M",
         "columns": [{"key": "a", "label": "A", "weight": 0},
                     {"key": "b", "label": "B", "weight": 0}],
         "rows": [{"a": "1", "b": "2"}]},
        {"type": "funnel", "title": "T", "primary_message": "M",
         "stages": [{"title": "a", "value": 0, "value_label": "0"},
                    {"title": "b", "value": 0, "value_label": "0"}]},
    ]
    for trial in degenerate:
        try:
            RENDERERS[trial["type"]](DeckBuilder(Path.cwd()), trial, 1)
        except Exception as exc:  # noqa: BLE001
            raise AssertionError(
                f"{trial['type']}: 退化した重み・値で生成が失敗しました: "
                f"{type(exc).__name__}: {exc}") from exc


def main() -> int:
    skill_root = Path(__file__).resolve().parents[1]
    sys.path.insert(0, str(skill_root / "runtime" / "python"))
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from slidekit import add_icon_list, inspect_content, logo_path_from_config, render_deck
    from slidekit.charts import add_native_chart
    from slidekit.icons import ICON_NAMES, add_icon
    from slidekit.preflight import KNOWN_TYPES
    from slidekit.theme import TYPE_BUSINESS
    from validate_pptx import validate

    config = {
        "python": {"executable": sys.executable},
        "paths": {
            "input_dir": "./build_slides/input", "work_dir": "./build_slides/work",
            "output_dir": "./build_slides/output",
        },
        "organization": {
            "department": "テスト部", "classification": "部外秘",
        },
        "deck": {"mode": "business"},
        "typography": {
            "headline_font": "Arial Unicode MS",
            "body_font": "Arial Unicode MS",
            "editorial_font": "Arial Unicode MS",
        },
        "branding": {"logo": {"enabled": True, "path": None}},
    }
    # 全renderer typeを最小構成で1枚ずつ網羅する。型ごとの回帰（例: フォント
    # サイズの取り違え、はみ出し、未対応フィールド）を機械的に検知するため。
    content = {
        "deck": {"mode": "business"},
        "slides": [
            {"id": "cover", "type": "cover", "title": "検証資料",
             "subtitle": "再現性テスト"},
            {"id": "comparison", "type": "comparison",
             "title": "比較して判断する", "density": "standard",
             "primary_message": "二つの観点を同時に確認する",
             "left": {"heading": "期待", "items": [
                 {"title": "効果", "body": "作業を短縮する"},
                 {"title": "品質", "body": "ばらつきを抑える"}]},
             "right": {"heading": "リスク", "items": [
                 {"title": "管理", "body": "無秩序な利用を防ぐ"},
                 {"title": "判断", "body": "基準を明確にする"}]}},
            {"id": "evidence", "type": "evidence_and_decision",
             "title": "根拠から判断する", "density": "standard",
             "primary_message": "検証してから広げる",
             "evidence_heading": "根拠", "evidence": [
                 {"title": "リスク", "body": "見極めが不十分"}],
             "decision_heading": "推奨", "decision_detail": "詳細"},
            {"id": "scope", "type": "scope_and_exclusions",
             "title": "範囲を絞る", "density": "dense",
             "primary_message": "対象を明確にする",
             "scope_heading": "範囲", "scope": [
                 {"label": "対象", "title": "業務A", "body": "説明"}],
             "exclusions_heading": "対象外", "exclusions": ["自動判断"]},
            {"id": "process", "type": "process_with_gates",
             "title": "段階的に進める", "density": "dense",
             "primary_message": "ゲートで判断する",
             "phases": [{"title": "準備", "label": "SETUP", "weight": 1,
                        "items": ["確定"]}],
             "gates": [{"title": "承認", "position": 1.0}]},
            {"id": "table", "type": "table_with_conclusion",
             "title": "表で評価する", "density": "dense",
             "primary_message": "必須項目を満たす",
             "columns": [{"key": "a", "label": "項目", "weight": 1}],
             "rows": [{"a": "作業時間削減", "_highlight": True}]},
            {"id": "table_insight", "type": "table_with_insight",
             "title": "表から複数の気づきを示す", "density": "dense",
             "primary_message": "気づきを踏まえた結論",
             "columns": [{"key": "a", "label": "項目", "weight": 1}],
             "rows": [{"a": "作業時間削減"}],
             "insights": ["気づき1"]},
            {"id": "org", "type": "org_layers", "title": "役割を分ける",
             "density": "standard", "primary_message": "責任を明確にする",
             "layers": [{"heading": "意思決定", "title": "経営層",
                        "body": "最終判断"}],
             "execution_heading": "実行", "execution": [
                 {"title": "部門A", "body": "利用"}]},
            {"id": "priority", "type": "priority_actions",
             "title": "優先度で対応する", "density": "standard",
             "primary_message": "重大リスクを優先する",
             "issues": [{"priority": "最優先", "title": "情報漏えい",
                        "body": "運用を徹底する"}],
             "actions": ["ログを記録する"]},
            {"id": "stage", "type": "stage_track", "title": "段階的に広げる",
             "density": "standard", "primary_message": "基準を満たせば進む",
             "stages": [{"label": "STEP1", "title": "PoC", "body": "検証"}]},
            {"id": "list", "type": "numbered_list", "title": "依頼事項",
             "density": "standard", "primary_message": "ご承認をお願いします",
             "message_position": "bottom",
             "items": [{"title": "予算承認", "body": "費用"}]},
            {"id": "divider", "type": "section_divider", "title": "体制編",
             "kicker": "SECTION 2"},
            {"id": "matrix", "type": "matrix", "title": "施策を整理する",
             "density": "standard", "primary_message": "効果が高いものから着手する",
             "x_axis": {"low": "易", "high": "難"},
             "y_axis": {"low": "小", "high": "大"},
             "cells": [
                 {"label": "優先", "title": "A", "body": "詳細", "emphasis": True},
                 {"label": "検討", "title": "B", "body": "詳細"},
                 {"label": "保留", "title": "C", "body": "詳細"},
                 {"label": "対象外", "title": "D", "body": "詳細"}]},
            {"id": "stat", "type": "stat_highlight", "title": "実績を示す",
             "density": "standard", "primary_message": "本格導入を検討する",
             "stat": {"value": "-42%", "label": "作業時間削減率"},
             "supporting": [{"value": "94%", "label": "継続意向"}]},
            {"id": "funnel", "type": "funnel", "title": "絞り込みを示す",
             "density": "standard", "primary_message": "資料請求後の転換率を上げる",
             "stages": [
                 {"title": "サイト訪問", "value": 12000, "value_label": "12,000"},
                 {"title": "資料請求", "value": 3400, "value_label": "3,400"},
                 {"title": "成約", "value": 210, "value_label": "210"}]},
            {"id": "cycle", "type": "cycle", "title": "繰り返しを示す",
             "density": "standard", "primary_message": "改善を継続する",
             "steps": [
                 {"label": "STEP1", "title": "Plan"}, {"label": "STEP2", "title": "Do"},
                 {"label": "STEP3", "title": "Check"}, {"label": "STEP4", "title": "Act"}]},
            {"id": "chart", "type": "chart_with_insight", "title": "結果を示す",
             "density": "standard", "primary_message": "全体像を確認する",
             "image": "chart.png", "insight_heading": "読み取れること",
             "insights": ["業務別の時間短縮率"]},
            {"id": "chart_native", "type": "chart_with_insight",
             "variant": "conclusion-led", "title": "実績を示す（ネイティブ）",
             "density": "standard", "primary_message": "実証期間で効果が拡大した",
             "chart": {"type": "column", "categories": ["準備", "実証", "評価"],
                      "series": [{"name": "削減率", "values": [12, 38, 42]}]}},
        ],
    }
    errors, warnings = inspect_content(content)
    assert not errors, errors
    assert set(s["type"] for s in content["slides"]) == KNOWN_TYPES, (
        "content が KNOWN_TYPES を網羅していません: "
        f"{KNOWN_TYPES - set(s['type'] for s in content['slides'])}")
    assert logo_path_from_config(config, Path.cwd()).is_file()

    # cover/section_divider以外の全rendererは spec["primary_message"] を直接
    # 参照するため、欠けたまま描画へ進むとKeyErrorで落ちる。preflightが必ず
    # errorで止める（＝警告どまりにしない）ことを型ごとに確認する。
    for slide in content["slides"]:
        if slide["type"] in {"cover", "section_divider"}:
            continue
        without_message = {k: v for k, v in slide.items() if k != "primary_message"}
        missing_errors, _ = inspect_content(
            {"deck": content["deck"], "slides": [without_message]})
        assert any("primary_message" in e for e in missing_errors), (
            f"{slide['type']}: primary_message欠落がerrorになっていません "
            "(描画時にKeyErrorで落ちます)")

    with tempfile.TemporaryDirectory(prefix="slidekit-test-") as temp:
        root = Path(temp)
        (root / ".slide-skill-config.yaml").write_text(
            yaml.safe_dump(config, allow_unicode=True), encoding="utf-8")
        input_dir = root / "build_slides" / "input"
        input_dir.mkdir(parents=True, exist_ok=True)
        from PIL import Image
        Image.new("RGB", (400, 240), (230, 236, 245)).save(input_dir / "chart.png")

        output, render_warnings = render_deck(content, root)
        assert not render_warnings, render_warnings
        presentation = Presentation(output)
        assert len(presentation.slides) == len(content["slides"])
        cover_text = [shape.text for shape in presentation.slides[0].shapes
                      if shape.has_text_frame]
        assert "部外秘" in cover_text and "テスト部" in cover_text

        issues, structure_warnings = validate(output)
        assert not issues, issues
        assert not structure_warnings, structure_warnings

        table_slide = presentation.slides[5]
        for shape in table_slide.shapes:
            if shape.has_table:
                for row in range(len(shape.table.rows)):
                    cell = shape.table.cell(row, 0)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            assert run.font.size.pt >= 11, (
                                f"表の文字が小さすぎます: {run.font.size.pt}pt "
                                f"'{run.text}'")

        icon_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        from pptx.util import Inches
        for index, name in enumerate(sorted(ICON_NAMES)):
            add_icon(icon_slide, Inches(0.5 + index), Inches(0.5), Inches(0.5), name)
        shape_count = len(icon_slide.shapes)
        assert shape_count >= len(ICON_NAMES) * 2, (
            f"アイコンの図形数が少なすぎます: {shape_count}")

        before = len(icon_slide.shapes)
        add_icon_list(icon_slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2),
                      ["一つ目の項目です", "二つ目の項目です", "三つ目の項目です"],
                      icon="check")
        new_text_boxes = [shape for shape in list(icon_slide.shapes)[before:]
                          if shape.has_text_frame and shape.text_frame.text.strip()]
        assert len(new_text_boxes) == 1, (
            f"add_icon_listの文章が複数shapeに分割されています: {len(new_text_boxes)}")
        assert len(new_text_boxes[0].text_frame.paragraphs) == 3

        comparison = presentation.slides[1]
        text_shapes = [shape for shape in comparison.shapes
                       if shape.has_text_frame and shape.text.strip()]
        assert len(text_shapes) <= 8, len(text_shapes)
        assert any(len(shape.text_frame.paragraphs) >= 2 for shape in text_shapes)

    # business以外のmodeでも同じcontentが構造的に破綻しない（はみ出し・エラー
    # なし）ことを確認する。denseは文字が小さく、large-roomは大きいため、
    # businessでは出ない回帰（はみ出し等）がmode依存で起きうる。
    import copy
    for mode in ("dense", "large-room"):
        mode_config = copy.deepcopy(config)
        mode_config["deck"]["mode"] = mode
        mode_content = copy.deepcopy(content)
        mode_content["deck"]["mode"] = mode
        with tempfile.TemporaryDirectory(prefix=f"slidekit-test-{mode}-") as temp:
            root = Path(temp)
            (root / ".slide-skill-config.yaml").write_text(
                yaml.safe_dump(mode_config, allow_unicode=True), encoding="utf-8")
            mode_input_dir = root / "build_slides" / "input"
            mode_input_dir.mkdir(parents=True, exist_ok=True)
            from PIL import Image as _Image
            _Image.new("RGB", (400, 240), (230, 236, 245)).save(
                mode_input_dir / "chart.png")

            mode_output, mode_render_warnings = render_deck(mode_content, root)
            assert not mode_render_warnings, (mode, mode_render_warnings)
            mode_issues, mode_warnings = validate(mode_output)
            assert not mode_issues, (mode, mode_issues)
            assert not mode_warnings, (mode, mode_warnings)

    _check_chart_stability(add_native_chart, TYPE_BUSINESS)
    _check_edge_case_stability(content)
    print("OK: slidekit self test")
    return 0


if __name__ == "__main__":
    sys.exit(main())
