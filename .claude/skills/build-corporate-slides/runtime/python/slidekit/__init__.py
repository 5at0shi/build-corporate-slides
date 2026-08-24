from pathlib import Path

from pptx import Presentation

from .components import (add_background_zone, add_card, add_cover, add_item_list,
                         add_icon_list, add_focus_panel, add_hairline, add_key_message,
                         add_numbered_row, add_panel, add_section_divider,
                         add_section_lead, add_slide_title)
from .theme import (LAYOUT, PALETTE, TYPE, TYPE_BUSINESS, TYPE_DENSE,
                    TYPE_LARGE_ROOM, TYPE_PRESENTATION, rgb,
                    typography_for)
from .typography import (add_paragraph_textbox, add_rich_textbox, add_textbox,
                         set_run, style_text_frame)


PLACEHOLDER_LOGO = Path(__file__).resolve().parents[3] / "assets" / "logo-placeholder.png"


def new_presentation(mode="business", fonts=None):
    presentation = Presentation()
    presentation.slide_width = LAYOUT.slide_width
    presentation.slide_height = LAYOUT.slide_height
    presentation._slidekit_typography = typography_for(mode, fonts=fonts)
    presentation._slidekit_mode = mode
    return presentation


def logo_path_from_config(config, workspace_root):
    """configのロゴ設定を解決する。enabled時の欠落は黙って無視しない。"""
    settings = config.get("branding", {}).get("logo", {})
    if not settings.get("enabled", False):
        return None
    configured = settings.get("path")
    if configured:
        candidate = Path(configured)
        if not candidate.is_absolute():
            candidate = Path(workspace_root) / candidate
    else:
        candidate = PLACEHOLDER_LOGO
    candidate = candidate.resolve()
    if not candidate.is_file():
        raise FileNotFoundError(f"ロゴ画像が見つかりません: {candidate}")
    return candidate


from .builder import DeckBuilder  # noqa: E402
from .charts import add_native_chart  # noqa: E402
from .config import load_workspace_config, workspace_paths  # noqa: E402
from .icons import ICON_NAMES, add_icon  # noqa: E402
from .layout import Region, content_region  # noqa: E402
from .images import add_image_contain  # noqa: E402
from .preflight import inspect_content, require_valid_content  # noqa: E402
from .renderers import RENDERERS, render_deck  # noqa: E402
from .tables import add_data_table  # noqa: E402


__all__ = [
    "LAYOUT", "PALETTE", "TYPE", "TYPE_BUSINESS", "TYPE_DENSE",
    "TYPE_LARGE_ROOM", "TYPE_PRESENTATION", "PLACEHOLDER_LOGO",
    "add_background_zone", "add_card", "add_cover", "add_item_list", "add_icon_list",
    "add_focus_panel",
    "add_panel", "add_section_divider",
    "add_hairline", "add_key_message", "add_numbered_row",
    "add_paragraph_textbox", "add_rich_textbox", "add_section_lead", "add_slide_title",
    "add_textbox", "new_presentation", "rgb", "set_run",
    "style_text_frame", "typography_for", "logo_path_from_config",
    "DeckBuilder", "Region", "content_region", "inspect_content",
    "load_workspace_config", "workspace_paths", "require_valid_content",
    "RENDERERS", "render_deck",
    "add_data_table", "add_image_contain",
    "add_native_chart", "add_icon", "ICON_NAMES",
]
